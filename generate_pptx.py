"""
day1company AI 전환 전략 — Part 1 (슬라이드 1~7)
slide_diagram_guide.yaml 기반 생성
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 디자인 시스템 상수 ──────────────────────────
# colors
C_DARK_BG    = RGBColor(0x1A, 0x1A, 0x1A)
C_WHITE_BG   = RGBColor(0xFF, 0xFF, 0xFF)
C_BOX_BG     = RGBColor(0xD9, 0xD9, 0xD9)
C_TEXT_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
C_SUB_DARK   = RGBColor(0xBB, 0xBB, 0xBB)
C_SUB_WHITE  = RGBColor(0x55, 0x55, 0x55)
C_CAPTION    = RGBColor(0x80, 0x80, 0x80)
C_ACCENT     = RGBColor(0xCC, 0x33, 0x33)

# font
FONT_FAMILY = "Noto Sans KR"
FONT_FAMILY_BLACK = "Noto Sans KR Black"

# canvas
CANVAS_W = Inches(10)
CANVAS_H = Inches(5.625)

# spacing_system.zones (white)
ZONE_TITLE_Y     = Inches(0.40)
ZONE_TITLE_H     = Inches(1.10)
ZONE_CONTENT_Y   = Inches(1.65)
ZONE_CONTENT_MAX_H = Inches(2.80)
ZONE_INSIGHT_Y   = Inches(4.85)
ZONE_INSIGHT_H   = Inches(0.35)

# spacing_system.dark_zones
DARK_MAIN_Y   = Inches(1.40)
DARK_MAIN_H   = Inches(1.80)
DARK_SUB_Y    = Inches(3.40)
DARK_SUB_H    = Inches(0.40)
DARK_META_Y   = Inches(4.60)
DARK_META_H   = Inches(0.35)

# margins
MARGIN_LEFT  = Inches(0.70)
CONTENT_W    = Inches(8.60)

# spacing scale
SP_XXS = Inches(0.04)
SP_XS  = Inches(0.08)
SP_SM  = Inches(0.10)
SP_MD  = Inches(0.15)
SP_LG  = Inches(0.20)
SP_XL  = Inches(0.30)

# equal_distribution presets
PRESETS = {
    2: {"w": Inches(4.20), "gap": Inches(0.20)},
    3: {"w": Inches(2.73), "gap": Inches(0.20)},
    4: {"w": Inches(2.05), "gap": Inches(0.10)},
    5: {"w": Inches(1.68), "gap": Inches(0.10)},
}

# flow_node
FLOW_NODE_H = Inches(0.38)
FLOW_ARROW_W = Inches(0.30)


# ── 텍스트 폭 추정 + 의미 단위 줄바꿈 함수 ────────
def estimate_text_width(text, font_size_pt, bold=False):
    """텍스트가 차지할 예상 폭(인치)을 근사 계산.
    한글: 글자당 약 font_size * 0.017"
    영문/숫자: 글자당 약 font_size * 0.009"
    Bold는 약 8% 추가"""
    width = 0
    for ch in text:
        if ord(ch) > 0x2E80:  # CJK 문자
            width += font_size_pt * 0.017
        else:
            width += font_size_pt * 0.009
    if bold:
        width *= 1.08
    return width


def semantic_line_break(text, box_width_inches, font_size, bold=False):
    """텍스트가 박스 폭을 초과하면, 의미 단위(어절/구두점)로 줄바꿈을 삽입.
    줄바꿈 우선순위: 쉼표(,) > 마침표(.) > 공백 > 기호(→, +, —)
    폰트 크기는 변경하지 않는다."""
    if '\n' in text:
        return text  # 이미 수동 줄바꿈이 있으면 그대로

    est = estimate_text_width(text, font_size, bold)
    if est <= box_width_inches:
        return text  # 한 줄에 들어가면 그대로

    # 줄바꿈 후보 위치 찾기 (우선순위 순)
    # 텍스트를 중간 지점 근처에서 끊기
    target_pos = len(text) // 2

    # 1차: 쉼표 근처에서 끊기
    best_pos = -1
    best_dist = len(text)
    for i, ch in enumerate(text):
        if ch in ',，' and i > len(text) * 0.3 and i < len(text) * 0.8:
            dist = abs(i - target_pos)
            if dist < best_dist:
                best_dist = dist
                best_pos = i + 1

    # 2차: 공백 근처에서 끊기
    if best_pos == -1:
        for i, ch in enumerate(text):
            if ch == ' ' and i > len(text) * 0.3 and i < len(text) * 0.8:
                dist = abs(i - target_pos)
                if dist < best_dist:
                    best_dist = dist
                    best_pos = i + 1

    # 3차: 기호(→, +, —) 근처에서 끊기
    if best_pos == -1:
        for i, ch in enumerate(text):
            if ch in '→+—' and i > len(text) * 0.25 and i < len(text) * 0.85:
                dist = abs(i - target_pos)
                if dist < best_dist:
                    best_dist = dist
                    best_pos = i

    if best_pos > 0:
        return text[:best_pos].rstrip() + '\n' + text[best_pos:].lstrip()

    return text  # 끊을 곳이 없으면 그대로


def fit_font_size(text, box_width_inches, max_font_size, bold=False, min_font_size=9):
    """텍스트가 박스 폭에 들어가는 최대 폰트 크기를 반환."""
    for fs in range(max_font_size, min_font_size - 1, -1):
        estimated = estimate_text_width(text, fs, bold)
        if estimated <= box_width_inches:
            return fs
    return min_font_size


# ── 헬퍼 함수 ──────────────────────────────────

def set_slide_bg(slide, color):
    """슬라이드 배경색 설정"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_run_style(run, font_size, bold, color, font_name):
    """run에 폰트 스타일 적용"""
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def _set_paragraph_text(tf, text, font_size, bold, color, font_name, align):
    """텍스트를 paragraph에 설정. \\n이 있으면 paragraph를 분리하여 각각에 동일 스타일 적용."""
    lines = text.split('\n') if '\n' in text else [text]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = line
        _set_run_style(run, font_size, bold, color, font_name)


def add_textbox(slide, left, top, width, height, text="",
                font_size=14, bold=False, color=C_TEXT_BLACK,
                font_name=FONT_FAMILY, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, fixed_size=False):
    """텍스트 박스 추가.
    fixed_size=True: 폰트 크기 고정 + 의미 단위 줄바꿈 (헤드라인, 중제목 등)
    fixed_size=False: 폰트 축소 허용 (본문, 보조 텍스트)"""
    box_w = width / 914400
    if text:
        if fixed_size:
            processed_text = semantic_line_break(text, box_w, font_size, bold)
            actual_fs = font_size
        else:
            processed_text = text
            actual_fs = fit_font_size(text, box_w, font_size, bold)
    else:
        processed_text = text
        actual_fs = font_size

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    _set_paragraph_text(tf, processed_text, actual_fs, bold, color, font_name, align)
    return txBox


def add_rect(slide, left, top, width, height,
             fill_color=None, border_color=None, border_width=Pt(0.5)):
    """사각형 도형 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text_to_shape(shape, text, font_size=12, bold=False,
                      color=C_TEXT_BLACK, font_name=FONT_FAMILY,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
                      fixed_size=False):
    """도형 내부에 텍스트 설정.
    fixed_size=True: 폰트 고정 + 의미 줄바꿈 (중제목 등)
    fixed_size=False: 폰트 축소 허용 (본문, 카드 내부 등)"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    box_w = shape.width / 914400
    usable_w = max(box_w - 0.15, 0.3)

    if text:
        if fixed_size:
            processed_text = semantic_line_break(text, usable_w, font_size, bold)
            actual_fs = font_size
        else:
            processed_text = text
            actual_fs = fit_font_size(text, usable_w, font_size, bold)
    else:
        processed_text = text
        actual_fs = font_size

    _set_paragraph_text(tf, processed_text, actual_fs, bold, color, font_name, align)
    return shape


def add_multiline_textbox(slide, left, top, width, height, lines,
                          anchor=MSO_ANCHOR.TOP, fixed_size=False):
    """여러 줄 텍스트 박스.
    fixed_size=True: 폰트 고정 + 의미 줄바꿈
    fixed_size=False: 폰트 축소 허용"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    box_w = width / 914400

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        text = line.get("text", "")
        fs = line.get("font_size", 14)
        bold = line.get("bold", False)

        if text:
            if fixed_size:
                processed = semantic_line_break(text, box_w, fs, bold)
                actual_fs = fs
            else:
                processed = text
                actual_fs = fit_font_size(text, box_w, fs, bold)
        else:
            processed = text
            actual_fs = fs

        p.text = processed
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = processed
        run.font.size = Pt(actual_fs)
        run.font.bold = bold
        run.font.color.rgb = line.get("color", C_TEXT_BLACK)
        run.font.name = line.get("font_name", FONT_FAMILY)
    return txBox


def add_insight_line(slide, text):
    """하단 인사이트 라인 (=> 접두어 + 18pt Bold 고정 + 중앙 정렬)
    18pt 고정. 폭 초과 시 의미 단위로 줄바꿈하여 2줄까지 허용.
    2줄이 되면 y좌표를 한 줄분 위로 올려서 공간 확보."""
    full_text = f"=> {text}"
    box_w = CONTENT_W / 914400  # 8.60"
    processed = semantic_line_break(full_text, box_w, 18, bold=True)

    # 2줄이면 높이를 0.60"로 확장하고 y를 위로
    is_two_lines = '\n' in processed
    insight_h = Inches(0.60) if is_two_lines else ZONE_INSIGHT_H
    insight_y = Inches(4.60) if is_two_lines else ZONE_INSIGHT_Y

    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, insight_y, CONTENT_W, insight_h
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    # processed에서 "=> " 접두어와 본문을 분리
    if processed.startswith("=> "):
        prefix = "=> "
        body = processed[3:]
    else:
        prefix = ""
        body = processed

    run1 = p.add_run()
    run1.text = prefix
    run1.font.size = Pt(18)
    run1.font.bold = True
    run1.font.color.rgb = C_TEXT_BLACK
    run1.font.name = FONT_FAMILY

    run2 = p.add_run()
    run2.text = body
    run2.font.size = Pt(18)
    run2.font.bold = True
    run2.font.color.rgb = C_TEXT_BLACK
    run2.font.name = FONT_FAMILY
    return txBox


def add_white_title(slide, text, chapter=None):
    """화이트 슬라이드 제목 (zones.title) — 20pt Bold, 줄바꿈 없이 한 줄
    chapter: 헤드라인 위에 작게 표시하는 챕터명 (10pt Bold)"""
    if chapter:
        add_textbox(slide, MARGIN_LEFT, ZONE_TITLE_Y, CONTENT_W, Inches(0.22),
                    text=chapter, font_size=10, bold=True,
                    color=C_SUB_WHITE, font_name=FONT_FAMILY)
        title_y = ZONE_TITLE_Y + Inches(0.22)
        title_h = ZONE_TITLE_H - Inches(0.22)
    else:
        title_y = ZONE_TITLE_Y
        title_h = ZONE_TITLE_H

    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, title_y, CONTENT_W, title_h
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0] if p.runs else p.add_run()
    if not p.runs:
        run.text = text
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = C_TEXT_BLACK
    run.font.name = FONT_FAMILY_BLACK


def card_x(n, idx):
    """N분할 배치에서 idx번째 카드의 x좌표"""
    preset = PRESETS[n]
    return MARGIN_LEFT + idx * (preset["w"] + preset["gap"])


# ── 슬라이드 생성 함수 ────────────────────────

def add_section_title(prs, title, subtitle=None):
    """중간 제목 슬라이드 (다크 배경) — Part 전환점. 줄바꿈 없이 원문 그대로."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK_BG)

    # 제목 — 줄바꿈/폰트 축소 없이 원문 그대로
    txBox = slide.shapes.add_textbox(MARGIN_LEFT, Inches(1.80), CONTENT_W, Inches(1.20))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0] if p.runs else p.add_run()
    if not p.runs:
        run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = C_TEXT_WHITE
    run.font.name = FONT_FAMILY_BLACK

    if subtitle:
        txBox2 = slide.shapes.add_textbox(MARGIN_LEFT, Inches(3.20), CONTENT_W, Inches(0.40))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.runs[0] if p2.runs else p2.add_run()
        if not p2.runs:
            run2.text = subtitle
        run2.font.size = Pt(14)
        run2.font.color.rgb = C_SUB_DARK
        run2.font.name = FONT_FAMILY


def slide_01_cover(prs):
    """슬라이드 1: 표지 (다크)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, C_DARK_BG)

    # 메인 텍스트
    add_multiline_textbox(
        slide, MARGIN_LEFT, DARK_MAIN_Y, CONTENT_W, DARK_MAIN_H,
        [
            {"text": "AI 기반 업무방식으로", "font_size": 34, "bold": True,
             "color": C_TEXT_WHITE, "font_name": FONT_FAMILY_BLACK, "align": PP_ALIGN.LEFT},
            {"text": "인당 생산성을 혁신한다", "font_size": 34, "bold": True,
             "color": C_TEXT_WHITE, "font_name": FONT_FAMILY_BLACK, "align": PP_ALIGN.LEFT},
        ]
    )

    # 부제
    add_textbox(
        slide, MARGIN_LEFT, DARK_SUB_Y, CONTENT_W, DARK_SUB_H,
        text="day1company 전사 AI 전환 전략", font_size=16,
        color=C_SUB_DARK, font_name=FONT_FAMILY
    )

    # 메타
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, DARK_META_Y, CONTENT_W, DARK_META_H
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    for part in [("발표자  ", 12, False), ("CAIO", 12, True),
                 ("    |    ", 12, False), ("2026-03-23", 12, False)]:
        run = p.add_run()
        run.text = part[0]
        run.font.size = Pt(part[1])
        run.font.bold = part[2]
        run.font.color.rgb = C_SUB_DARK
        run.font.name = FONT_FAMILY


def slide_02_global_cases(prs):
    """슬라이드 2: 글로벌 기업 사례 — 5열 카드 (재디자인)
    개선: 핵심 수치를 다크 배경 강조 영역으로 분리, 기업명/업종을 상단 헤더로,
    결과를 하단 독립 영역으로 시각적 위계를 명확하게 구성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)

    add_white_title(slide, "글로벌 선도 기업들은 AI로 비용 구조 자체를 전환 중", chapter="Part 1 — 비용 구조가 바뀌고 있다")

    companies = [
        {"name": "Klarna",   "industry": "핀테크",   "metric": "-47%",    "detail": "5,527→2,907명", "result": "인당 매출 73%↑"},
        {"name": "Block",    "industry": "핀테크",   "metric": "-40%",    "detail": "10K→6K명",       "result": "연 $600M 절감"},
        {"name": "Amazon",   "industry": "클라우드", "metric": "-30K명",  "detail": "Project Dawn",    "result": "$2.1~3.6B 절감"},
        {"name": "Shopify",  "industry": "이커머스", "metric": "-30%",    "detail": "11.6K→8.1K명",   "result": "인당 매출 127%↑"},
        {"name": "Duolingo", "industry": "에듀테크", "metric": "4~5x",    "detail": "외주→AI 전환",    "result": "해고 0명"},
    ]

    n = 5
    card_w = PRESETS[n]["w"]
    card_y = ZONE_CONTENT_Y

    # 3개 영역: 헤더(기업명+업종) / 핵심 수치(다크) / 결과(그레이)
    header_h = Inches(0.55)
    metric_h = Inches(1.10)
    result_h = Inches(0.65)
    total_h = header_h + metric_h + result_h  # 2.30"

    for i, co in enumerate(companies):
        x = card_x(n, i)

        # 상단 헤더 영역 (흰 배경 + 테두리)
        add_rect(slide, x, card_y, card_w, header_h,
                 fill_color=C_WHITE_BG, border_color=C_DARK_BG)
        add_textbox(slide, x, card_y + SP_XXS, card_w, Inches(0.30),
                    text=co["name"], font_size=15, bold=True,
                    color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x, card_y + Inches(0.32), card_w, Inches(0.20),
                    text=co["industry"], font_size=9,
                    color=C_CAPTION, align=PP_ALIGN.CENTER)

        # 중앙 핵심 수치 영역 (다크 배경 — 시각적 무게 집중)
        metric_y = card_y + header_h
        add_rect(slide, x, metric_y, card_w, metric_h,
                 fill_color=C_DARK_BG, border_color=C_DARK_BG)
        add_textbox(slide, x, metric_y + Inches(0.15), card_w, Inches(0.55),
                    text=co["metric"], font_size=28, bold=True,
                    color=C_TEXT_WHITE, font_name=FONT_FAMILY_BLACK,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x, metric_y + Inches(0.72), card_w, Inches(0.25),
                    text=co["detail"], font_size=9,
                    color=C_SUB_DARK, align=PP_ALIGN.CENTER)

        # 하단 결과 영역 (그레이 배경)
        result_y = metric_y + metric_h
        add_rect(slide, x, result_y, card_w, result_h,
                 fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_textbox(slide, x, result_y + Inches(0.12), card_w, Inches(0.40),
                    text=co["result"], font_size=12, bold=True,
                    color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)

    add_insight_line(slide, "AI '도입'이 아닌 비용 구조 자체의 전환")


def _add_flow_nodes(slide, nodes, start_x, y, node_w, node_h, arrow_w,
                    node_fill, node_border, text_color, strikethrough_idx=None):
    """흐름도 노드 배치 헬퍼"""
    x = start_x
    for i, text in enumerate(nodes):
        if i > 0:
            # 화살표
            add_textbox(slide, x, y, arrow_w, node_h,
                        text="→", font_size=14, bold=True,
                        color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)
            x += arrow_w

        shape = add_rect(slide, x, y, node_w, node_h,
                         fill_color=node_fill, border_color=node_border)
        fs = 10
        b = False
        c = text_color
        if i == len(nodes) - 1:
            b = True
        add_text_to_shape(shape, text, font_size=fs, bold=b, color=c)
        x += node_w
    return x


def slide_03_new_normal(prs):
    """슬라이드 3: New Normal — 2분할 Before/After (재디자인)
    개선: 각 패턴을 Before→After 흐름으로 재구성.
    흐름도 노드를 크게, 취소선 노드를 시각적으로 구분.
    사례를 흐름 아래에 컴팩트하게 배치."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)

    add_white_title(slide, "채용과 운영의 New Normal: 사람을 늘리지 않고 AI로 처리", chapter="Part 1 — 비용 구조가 바뀌고 있다")

    n = 2
    card_w = PRESETS[n]["w"]
    card_y = ZONE_CONTENT_Y

    panels = [
        {
            "label": "패턴 1",
            "title": "퇴사·결원 시 → 충원 대신 AI 효율화",
            "flow": [
                ("퇴사·결원", C_BOX_BG, C_TEXT_BLACK, False),
                ("충원", C_BOX_BG, C_SUB_WHITE, True),  # 취소 — #555555
                ("AI 효율화", C_DARK_BG, C_TEXT_WHITE, False),  # 강조
            ],
            "cases": [
                ("Klarna", "47% 감소, 인당 매출 73%↑"),
                ("Block", "40% 감축, 연 $600M 절감"),
                ("Duolingo", "외주→AI, 생산량 4~5배"),
            ]
        },
        {
            "label": "패턴 2",
            "title": "인력 필요 시 → 채용 대신 AI 생산성 증대",
            "flow": [
                ("인력 필요", C_BOX_BG, C_TEXT_BLACK, False),
                ("채용", C_BOX_BG, C_SUB_WHITE, True),  # 취소 — #555555
                ("AI 생산성↑", C_DARK_BG, C_TEXT_WHITE, False),  # 강조
            ],
            "cases": [
                ("Shopify", "AI 불가 소명 의무, 인당 매출 127%↑"),
                ("Amazon", "30K 감축 + AI에 $125B 투자"),
                ("Klarna", "마케팅팀 50%↓, 캠페인 수↑"),
            ]
        },
    ]

    for i, panel in enumerate(panels):
        x = card_x(n, i)

        # 라벨 (다크 배지)
        badge = add_rect(slide, x, card_y, Inches(0.80), Inches(0.30),
                         fill_color=C_DARK_BG, border_color=C_DARK_BG)
        add_text_to_shape(badge, panel["label"], font_size=11, bold=True, color=C_TEXT_WHITE)

        # 제목
        add_textbox(slide, x + Inches(0.90), card_y, card_w - Inches(0.90), Inches(0.30),
                    text=panel["title"], font_size=13, bold=True,
                    color=C_TEXT_BLACK, align=PP_ALIGN.LEFT)

        # 흐름도 — 넓은 노드
        flow_y = card_y + Inches(0.50)
        node_w = Inches(1.10)
        node_h = Inches(0.48)
        arrow_w = Inches(0.25)
        fx = x

        for j, (text, fill, tc, is_cancelled) in enumerate(panel["flow"]):
            if j > 0:
                add_textbox(slide, fx, flow_y, arrow_w, node_h,
                            text="→", font_size=16, bold=True,
                            color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)
                fx += arrow_w

            box = add_rect(slide, fx, flow_y, node_w, node_h,
                           fill_color=fill, border_color=C_DARK_BG)
            fs = 12 if not is_cancelled else 11
            add_text_to_shape(box, text, font_size=fs, bold=(not is_cancelled), color=tc)

            # 취소 표시 — 가로선
            if is_cancelled:
                line_y = flow_y + node_h / 2
                add_rect(slide, fx + SP_SM, line_y, node_w - SP_SM * 2, Pt(1.5),
                         fill_color=C_ACCENT)
            fx += node_w

        # 사례 — 기업명 Bold + 수치
        case_y = card_y + Inches(1.20)
        for company, detail in panel["cases"]:
            txBox = slide.shapes.add_textbox(x + SP_SM, case_y, card_w - SP_SM * 2, Inches(0.22))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r1 = p.add_run()
            r1.text = f"{company}  "
            r1.font.size = Pt(10)
            r1.font.bold = True
            r1.font.color.rgb = C_TEXT_BLACK
            r1.font.name = FONT_FAMILY
            r2 = p.add_run()
            r2.text = detail
            r2.font.size = Pt(10)
            r2.font.color.rgb = C_TEXT_BLACK
            r2.font.name = FONT_FAMILY
            case_y += Inches(0.25)

    # 하단 결과 바
    bar = add_rect(slide, MARGIN_LEFT, Inches(4.35), CONTENT_W, Inches(0.40),
                   fill_color=C_DARK_BG, border_color=C_DARK_BG)
    add_text_to_shape(bar,
        "결과: 남은 인력 → 더 높은 급여 (Klarna 60%↑) + 더 높은 인당 생산성 + 더 높은 가치의 업무",
        font_size=12, bold=True, color=C_TEXT_WHITE)


def slide_04_agenda(prs):
    """슬라이드 4: 목차 (화이트, 7행 진행 흐름도)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)

    add_white_title(slide, "오늘의 논의: AI 전환 전략의 구조, 실행, 의사결정", chapter="Part 1 — 비용 구조가 바뀌고 있다")

    parts = [
        ("Part 1", "AI 시대, 비용 구조가 바뀌고 있다", "1~7"),
        ("Part 2", "AI 전환 전략 프레임워크", "8~12"),
        ("Part 3", "Layer 1 — 첫 프로젝트 선정과 2개월 실행", "13~18"),
        ("Part 4", "Layer 2 — 전사 AI 활용 제도 설계", "19~23"),
        ("Part 5", "통합 로드맵과 기대 효과", "24~26"),
        ("Part 6", "예상 리스크와 대응 설계", "27~29"),
        ("Part 7", "의사결정 요청 + Q&A", "30~32"),
    ]

    row_h = Inches(0.40)
    row_gap = Inches(0.06)
    start_y = ZONE_CONTENT_Y

    for i, (part, title, pages) in enumerate(parts):
        y = start_y + i * (row_h + row_gap)

        # Part 번호 (다크 박스)
        part_box = add_rect(slide, MARGIN_LEFT, y, Inches(0.80), row_h,
                            fill_color=C_DARK_BG, border_color=C_DARK_BG)
        add_text_to_shape(part_box, part, font_size=12, bold=True,
                          color=C_TEXT_WHITE)

        # 제목 영역 (그레이 배경)
        title_box = add_rect(slide, Inches(1.58), y, Inches(6.80), row_h,
                             fill_color=C_BOX_BG, border_color=None)
        # 제목 텍스트
        add_textbox(slide, Inches(1.70), y, Inches(5.60), row_h,
                    text=title, font_size=14, color=C_TEXT_BLACK,
                    align=PP_ALIGN.LEFT)

        # 페이지 번호
        add_textbox(slide, Inches(7.60), y, Inches(0.70), row_h,
                    text=pages, font_size=10, color=C_SUB_WHITE,
                    align=PP_ALIGN.CENTER)

    add_insight_line(slide, "L1 첫 프로젝트 선정 + L2 제도 도입 승인")


def slide_05_disruption(prs):
    """슬라이드 5: 우리가 기존 업계에 했던 것을 누군가가 우리에게 (화이트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)

    add_white_title(slide, "우리가 성장한 방식으로 추월 당할 수 있다는 위기감", chapter="Part 1 — 비용 구조가 바뀌고 있다")

    # 중앙 타임라인
    line_y = Inches(2.75)
    add_rect(slide, MARGIN_LEFT, line_y, CONTENT_W, Pt(1.5),
             fill_color=C_DARK_BG)

    # "지금" 라벨
    now_box = add_rect(slide, Inches(4.50), Inches(2.55), Inches(1.00), Inches(0.40),
                       fill_color=C_DARK_BG, border_color=C_DARK_BG)
    add_text_to_shape(now_box, "지금", font_size=12, bold=True, color=C_TEXT_WHITE)

    # 좌측: 2010년대
    add_textbox(slide, MARGIN_LEFT, Inches(2.95), Inches(3.50), Inches(0.30),
                text="2010년대", font_size=12, color=C_SUB_WHITE)

    # day1company (추월하는 쪽)
    left_top = add_rect(slide, Inches(0.90), ZONE_CONTENT_Y, Inches(2.80), Inches(0.45),
                        fill_color=C_DARK_BG, border_color=C_DARK_BG)
    add_text_to_shape(left_top, "day1company (Digital-native)",
                      font_size=12, bold=True, color=C_TEXT_WHITE)

    add_textbox(slide, Inches(1.20), Inches(2.15), Inches(2.20), Inches(0.30),
                text="── 추월 ──▶", font_size=12, bold=True,
                color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)

    left_bot = add_rect(slide, Inches(0.90), Inches(2.45), Inches(2.80), Inches(0.40),
                        fill_color=C_BOX_BG, border_color=C_DARK_BG)
    add_text_to_shape(left_bot, "기존 교육업계 (오프라인)",
                      font_size=12, color=C_TEXT_BLACK)

    # 우측: 202X년
    add_textbox(slide, Inches(5.80), Inches(2.95), Inches(3.50), Inches(0.30),
                text="202X년", font_size=12, color=C_SUB_WHITE)

    right_top = add_rect(slide, Inches(6.30), ZONE_CONTENT_Y, Inches(2.80), Inches(0.45),
                         fill_color=C_DARK_BG, border_color=C_DARK_BG)
    add_text_to_shape(right_top, "AI-Native 신규 진입자",
                      font_size=12, bold=True, color=C_TEXT_WHITE)

    add_textbox(slide, Inches(6.60), Inches(2.15), Inches(2.20), Inches(0.30),
                text="── 추월 ──▶", font_size=12, bold=True,
                color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)

    right_bot = add_rect(slide, Inches(6.30), Inches(2.45), Inches(2.80), Inches(0.40),
                         fill_color=C_BOX_BG, border_color=C_DARK_BG)
    add_text_to_shape(right_bot, "day1company (현재)",
                      font_size=12, color=C_TEXT_BLACK)

    # 하단 설명 박스
    desc_box = add_rect(slide, MARGIN_LEFT, Inches(3.35), CONTENT_W, Inches(1.20),
                        fill_color=C_BOX_BG, border_color=C_DARK_BG)
    add_multiline_textbox(
        slide, Inches(0.90), Inches(3.45), Inches(8.20), Inches(1.00),
        [
            {"text": "과거: day1company는 Digital-native 조직으로서 기존 교육업계를 추월했다",
             "font_size": 12, "bold": False, "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT},
            {"text": "현재: AI-Native 조직들이 처음부터 AI를 전제로 설계 — 10명이 100명분의 일을 처리",
             "font_size": 12, "bold": False, "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT},
            {"text": "결론: Digital-native에서 AI-native로의 전환은 경쟁 우위를 유지하기 위한 필수 조건",
             "font_size": 12, "bold": True, "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT},
        ]
    )

    add_insight_line(slide, "업무 방식 전환이 경쟁력의 유일한 방법")


def slide_06_current_state(prs):
    """슬라이드 6: day1company 현황 (재디자인)
    개선: 좌측 가로 막대를 실제 비례 길이로, 시간 수치를 막대 끝에 배치.
    우측 매트릭스를 축 라벨 추가하여 의미 명확화. 하단 요약을 전폭 배너로."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)

    add_white_title(slide, "매출이 성장하면 인건비도 비례하여 증가하는 구조", chapter="Part 1 — 비용 구조가 바뀌고 있다")

    # ── 좌측: 월간 반복 업무 시간 ──
    left_x = MARGIN_LEFT
    left_w = Inches(4.40)
    add_textbox(slide, left_x, ZONE_CONTENT_Y, left_w, Inches(0.25),
                text="월간 반복 업무 시간", font_size=14, bold=True,
                color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK)

    # 시간 기준 최대값 32h → 최대 막대 폭
    max_bar_w = Inches(3.40)
    tasks = [
        ("매출 정산", "~32h", 32),
        ("비용 정산", "~24h", 24),
        ("보고·집계", "수시간~수일", 16),
        ("데이터 가공", "공통", 12),
    ]

    bar_y = Inches(2.00)
    bar_h = Inches(0.42)
    bar_gap = Inches(0.08)
    label_w = Inches(1.00)

    for i, (name, time_str, hours) in enumerate(tasks):
        y = bar_y + i * (bar_h + bar_gap)
        bar_w_actual = max_bar_w * (hours / 32)

        # 라벨 (좌측)
        add_textbox(slide, left_x, y, label_w, bar_h,
                    text=name, font_size=11, bold=True,
                    color=C_TEXT_BLACK, align=PP_ALIGN.LEFT)

        # 막대 (비례 길이)
        bar_x = left_x + label_w + SP_SM
        fill = C_DARK_BG if i == 0 else C_BOX_BG
        tc = C_TEXT_WHITE if i == 0 else C_TEXT_BLACK
        bar = add_rect(slide, bar_x, y, bar_w_actual, bar_h,
                       fill_color=fill, border_color=C_DARK_BG)

        # 시간 (막대 내부 우측)
        add_textbox(slide, bar_x + bar_w_actual - Inches(0.90), y, Inches(0.85), bar_h,
                    text=time_str, font_size=11, bold=True,
                    color=tc, align=PP_ALIGN.RIGHT)

    # ── 우측: 곱 증가 매트릭스 ──
    right_x = Inches(5.50)
    right_w = Inches(4.00)
    add_textbox(slide, right_x, ZONE_CONTENT_Y, right_w, Inches(0.25),
                text="성장 시 업무량 곱 증가", font_size=14, bold=True,
                color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK)

    # 축 라벨
    add_textbox(slide, right_x, Inches(2.00), right_w, Inches(0.22),
                text="사업부문 7 × 상품유형 6 = 42 조합", font_size=11,
                bold=True, color=C_TEXT_BLACK)

    # 매트릭스 (7열 × 6행) — 첫 행/첫 열을 다크로 축 표시
    cell_w = Inches(0.44)
    cell_h = Inches(0.28)
    cell_gap = Inches(0.02)
    grid_x = right_x
    grid_y = Inches(2.30)

    for row in range(6):
        for col in range(7):
            cx = grid_x + col * (cell_w + cell_gap)
            cy = grid_y + row * (cell_h + cell_gap)
            # 첫 행 또는 첫 열을 다크로 — 축 느낌
            if row == 0 or col == 0:
                fill = C_DARK_BG
            else:
                fill = C_BOX_BG
            add_rect(slide, cx, cy, cell_w, cell_h,
                     fill_color=fill, border_color=C_DARK_BG, border_width=Pt(0.25))

    add_textbox(slide, right_x, Inches(4.15), right_w, Inches(0.20),
                text="각 칸 = 정산·보고 1건 → 사업 확장 시 칸이 늘어남", font_size=9,
                color=C_CAPTION)

    # ── 하단 전폭 요약 ──
    summary = add_rect(slide, MARGIN_LEFT, Inches(4.35), CONTENT_W, Inches(0.40),
                       fill_color=C_DARK_BG, border_color=C_DARK_BG)
    add_text_to_shape(summary,
        "7개 사업부문 | 인당 매출 2.34억 원 (25년 연말 기준) | 매출↑ = 인건비↑ 구조",
        font_size=12, bold=True, color=C_TEXT_WHITE)

    add_insight_line(slide, "매출↑ = 인건비↑ 구조")


def slide_07_before_after(prs):
    """슬라이드 7: AI 전환이 해결하는 문제와 목표 (화이트, Before/After)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)

    add_white_title(slide, "AI 전환은 '성장 → 인력 추가' 고리를 끊는다", chapter="Part 1 — 비용 구조가 바뀌고 있다")

    # Before 라벨
    add_textbox(slide, MARGIN_LEFT, ZONE_CONTENT_Y, Inches(0.80), Inches(0.25),
                text="Before", font_size=12, bold=True, color=C_TEXT_BLACK)

    # Before 흐름
    before_nodes = ["사업 성장", "업무량↑", "인력 채용↑", "인건비↑"]
    node_w = Inches(1.50)
    node_h = Inches(0.42)
    before_y = Inches(1.95)

    x = Inches(1.00)
    for i, text in enumerate(before_nodes):
        if i > 0:
            add_textbox(slide, x, before_y, Inches(0.50), node_h,
                        text="→", font_size=14, bold=True,
                        color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)
            x += Inches(0.50)

        box = add_rect(slide, x, before_y, node_w, node_h,
                       fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_text_to_shape(box, text, font_size=12, color=C_TEXT_BLACK)
        x += node_w

    # After 라벨
    add_textbox(slide, MARGIN_LEFT, Inches(2.65), Inches(0.80), Inches(0.25),
                text="After", font_size=12, bold=True, color=C_TEXT_BLACK)

    # After 흐름
    after_nodes = ["사업 성장", "업무량↑", "AI 처리", "인력 동결", "인당 생산성↑"]
    after_y = Inches(2.95)
    node_w2 = Inches(1.25)

    x = Inches(1.00)
    for i, text in enumerate(after_nodes):
        if i > 0:
            add_textbox(slide, x, after_y, Inches(0.30), node_h,
                        text="→", font_size=14, bold=True,
                        color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)
            x += Inches(0.30)

        fill = C_DARK_BG if text == "AI 처리" else C_BOX_BG
        tc = C_TEXT_WHITE if text == "AI 처리" else C_TEXT_BLACK
        b = text == "AI 처리"
        box = add_rect(slide, x, after_y, node_w2, node_h,
                       fill_color=fill, border_color=C_DARK_BG)
        add_text_to_shape(box, text, font_size=11, bold=b, color=tc)
        x += node_w2

    # 우측 목표 수치 박스
    target_box = add_rect(slide, Inches(5.80), Inches(3.50), Inches(3.50), Inches(1.00),
                          fill_color=C_BOX_BG, border_color=C_DARK_BG)
    add_textbox(slide, Inches(5.80), Inches(3.50), Inches(3.50), Inches(0.50),
                text="인당 매출 28%↑", font_size=22, bold=True,
                color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.80), Inches(4.00), Inches(3.50), Inches(0.30),
                text="2.34억 → 3억 원", font_size=14,
                color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.80), Inches(4.30), Inches(3.50), Inches(0.20),
                text="참고: Klarna 73%↑ / Shopify 127%↑ 대비 보수적 목표",
                font_size=8, color=C_CAPTION, align=PP_ALIGN.CENTER)

    # 좌측 달성 경로
    add_multiline_textbox(
        slide, MARGIN_LEFT, Inches(3.65), Inches(4.80), Inches(0.35),
        [
            {"text": "달성 경로: Layer 1 (반복 업무 자동화) + Layer 2 (전사 AI 활용 제도)",
             "font_size": 12, "bold": False, "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT},
        ]
    )

    add_insight_line(slide, "보수적 목표 — 실행 방법을 지금부터 설명")


# ══════════════════════════════════════════════════
# Part 2: AI 전환 전략 프레임워크 (슬라이드 8~12)
# ══════════════════════════════════════════════════

def slide_08_two_layers(prs):
    """슬라이드 8: AI 전환의 두 가지 Layer (화이트, 2분할)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "전사 AI 전환은 Top-down과 Bottom-up, 두 가지 Layer로 진행", chapter="Part 2 — AI 전환 전략 프레임워크")

    n = 2
    card_w = PRESETS[n]["w"]
    card_h = Inches(2.80)
    card_y = ZONE_CONTENT_Y

    layers = [
        {
            "title": "Layer 1", "subtitle": "구조화 자동화 — Palantir Approach",
            "items": [
                ("예시", "매출 정산, 강사료 정산, 광고비 정산, 손익 보고"),
                ("출력 성격", "결정적 — 같은 입력이면 항상 같은 결과"),
                ("추진 주체", "CAIO 주도 (Top-down)"),
            ]
        },
        {
            "title": "Layer 2", "subtitle": "AI 도구 활용 효율화",
            "items": [
                ("예시", "상세페이지 생성, 고객 응대, 회의록, 보고서 초안, 번역"),
                ("출력 성격", "확률적 — 매번 다른 결과 가능"),
                ("추진 주체", "각 부서 자율 (Bottom-up)"),
            ]
        },
    ]

    for i, layer in enumerate(layers):
        x = card_x(n, i)
        add_rect(slide, x, card_y, card_w, card_h,
                 fill_color=C_BOX_BG, border_color=C_DARK_BG)

        # 타이틀
        add_textbox(slide, x + SP_LG, card_y + Inches(0.12), Inches(2.00), Inches(0.35),
                    text=layer["title"], font_size=16, bold=True,
                    color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK)
        add_textbox(slide, x + SP_LG, card_y + Inches(0.45), card_w - SP_LG * 2, Inches(0.25),
                    text=layer["subtitle"], font_size=12, color=C_TEXT_BLACK)

        # 항목
        item_y = card_y + Inches(0.85)
        lines = []
        for label, value in layer["items"]:
            lines.append({"text": label, "font_size": 10, "bold": True,
                          "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT})
            b = label == "추진 주체"
            lines.append({"text": value, "font_size": 12, "bold": b,
                          "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT})
        add_multiline_textbox(slide, x + SP_LG, item_y,
                              card_w - SP_LG * 2, Inches(1.90), lines)

    # 중앙 양방향 화살표 — 카드 아래 중앙에 배치
    arrow_y = card_y + card_h + SP_XXS  # 카드 바로 아래
    arrow_box = add_rect(slide, Inches(4.20), arrow_y, Inches(1.20), Inches(0.30),
                         fill_color=C_DARK_BG, border_color=C_DARK_BG)
    add_text_to_shape(arrow_box, "⟷ 상호보완", font_size=11, bold=True, color=C_TEXT_WHITE)

    add_insight_line(slide, "Top-down + Bottom-up 병행 구조")


def slide_09_synergy(prs):
    """슬라이드 9: 상호보완 구조 — 사이클 다이어그램 (재디자인)
    개선: 텍스트 나열 → 좌측 사이클 다이어그램 + 우측 핵심 설명 3건"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "Layer 2의 부서별 성과가 Layer 1에 편입될 때 지속 가능한 절감 가능", chapter="Part 2 — AI 전환 전략 프레임워크")

    # ── 좌측: 사이클 다이어그램 (4노드 직사각형 배치) ──
    # 명확한 좌표: 상/하 노드는 중앙 정렬, 좌/우 노드는 양옆
    # 화살표는 노드 사이 정확한 중간 위치

    node_w = Inches(1.60)
    node_h = Inches(0.42)
    arrow_size = Inches(0.30)
    gap_h = Inches(0.30)  # 상-좌/우 간 세로 간격
    gap_w = Inches(0.30)  # 좌-상/하 간 가로 간격

    # 기준점: 좌측 영역 중앙 (0.70"~5.50" 사이)
    area_cx = Inches(3.10)  # 좌측 영역 중앙 x
    area_cy = Inches(3.05)  # 콘텐츠 영역 중앙 y

    # 4노드 좌표 계산
    top_x = area_cx - node_w / 2
    top_y = area_cy - node_h - gap_h - arrow_size / 2

    bot_x = area_cx - node_w / 2
    bot_y = area_cy + gap_h + arrow_size / 2

    left_x = area_cx - node_w - gap_w - arrow_size / 2
    left_x = max(MARGIN_LEFT, left_x)  # 마진 보정
    left_y = area_cy - node_h / 2

    right_x = area_cx + gap_w + arrow_size / 2
    right_y = area_cy - node_h / 2

    nodes = [
        (top_x, top_y, "L1: 전사 구조 구축", C_DARK_BG, C_TEXT_WHITE),
        (right_x, right_y, "부서별 AI 활용", C_BOX_BG, C_TEXT_BLACK),
        (bot_x, bot_y, "L2: 부서 성과 발생", C_BOX_BG, C_TEXT_BLACK),
        (left_x, left_y, "성과 → L1 편입", C_DARK_BG, C_TEXT_WHITE),
    ]

    for nx, ny, text, fill, tc in nodes:
        box = add_rect(slide, nx, ny, node_w, node_h,
                       fill_color=fill, border_color=C_DARK_BG)
        add_text_to_shape(box, text, font_size=10, bold=True, color=tc)

    # 시계방향 화살표 — 노드 사이 정확한 중간
    # 상→우: 상 노드 우하단과 우 노드 좌상단 사이
    ar_tr_x = top_x + node_w + (right_x - top_x - node_w) / 2 - arrow_size / 2
    ar_tr_y = top_y + node_h + (right_y - top_y - node_h) / 2 - arrow_size / 2
    add_textbox(slide, ar_tr_x, ar_tr_y, arrow_size, arrow_size,
                text="↘", font_size=16, bold=True, color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)

    # 우→하: 우 노드 좌하단과 하 노드 우상단 사이
    ar_rb_x = right_x + (bot_x + node_w - right_x - node_w) / 2
    ar_rb_y = right_y + node_h + (bot_y - right_y - node_h) / 2 - arrow_size / 2
    add_textbox(slide, ar_rb_x, ar_rb_y, arrow_size, arrow_size,
                text="↙", font_size=16, bold=True, color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)

    # 하→좌: 하 노드 좌상단과 좌 노드 우하단 사이
    ar_bl_x = left_x + node_w + (bot_x - left_x - node_w) / 2 - arrow_size / 2
    ar_bl_y = left_y + node_h + (bot_y - left_y - node_h) / 2 - arrow_size / 2
    add_textbox(slide, ar_bl_x, ar_bl_y, arrow_size, arrow_size,
                text="↖", font_size=16, bold=True, color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)

    # 좌→상: 좌 노드 우상단과 상 노드 좌하단 사이
    ar_lt_x = left_x + node_w + (top_x - left_x - node_w) / 2 - arrow_size / 2
    ar_lt_y = top_y + node_h + (left_y - top_y - node_h) / 2 - arrow_size / 2
    add_textbox(slide, ar_lt_x, ar_lt_y, arrow_size, arrow_size,
                text="↗", font_size=16, bold=True, color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)

    # ── 우측: 핵심 설명 3건 ──
    rx = Inches(5.80)
    rw = Inches(3.50)
    ry_start = ZONE_CONTENT_Y + Inches(0.10)

    explanations = [
        ("①", "Top-down + Bottom-up = 전사 변화",
         "CAIO가 구조를 만들고,\n부서가 그 위에서 AI 활용"),
        ("②", "L2 → L1 편입 경로",
         "부서 AI 성과가 반복·검증되면\nL1 워크플로우에 편입"),
        ("③", "표준화 = 지식 보존",
         "사람 의존 → 시스템 전환\n인력 변동에도 유지"),
    ]

    ey = ry_start
    for num, title, desc in explanations:
        # 번호
        num_box = add_rect(slide, rx, ey, Inches(0.28), Inches(0.28),
                           fill_color=C_DARK_BG, border_color=C_DARK_BG)
        add_text_to_shape(num_box, num, font_size=11, bold=True, color=C_TEXT_WHITE)

        add_textbox(slide, rx + Inches(0.35), ey, rw - Inches(0.35), Inches(0.25),
                    text=title, font_size=12, bold=True, color=C_TEXT_BLACK)

        add_multiline_textbox(
            slide, rx + Inches(0.35), ey + Inches(0.28), rw - Inches(0.35), Inches(0.45),
            [{"text": line, "font_size": 10, "bold": False, "color": C_SUB_WHITE, "align": PP_ALIGN.LEFT}
             for line in desc.split("\n")]
        )
        ey += Inches(0.85)

    # 하단 요약
    summary = add_rect(slide, MARGIN_LEFT, Inches(4.40), CONTENT_W, Inches(0.35),
                       fill_color=C_BOX_BG, border_color=C_DARK_BG)
    add_text_to_shape(summary,
        "핵심: L2 부서 성과 → 반복·검증 → L1 시스템 편입 = 지속 가능한 절감",
        font_size=11, bold=True, color=C_TEXT_BLACK)


def slide_10_palantir(prs):
    """슬라이드 10: Layer 1 참고 사례 (화이트, 4열 카드)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "Layer 1의 데이터 구조화 + 자동화 방식은 글로벌 기업에서 검증", chapter="Part 2 — AI 전환 전략 프레임워크")

    n = 4
    card_w = PRESETS[n]["w"]
    card_h = Inches(2.50)
    card_y = ZONE_CONTENT_Y

    types = [
        {"label": "유형 1", "title_l1": "수작업 자동화", "title_l2": "→ 인건비 절감",
         "cases": ["AML 비용 90%↓", "투자 리드 100배↑", "간호사 편성 자동화"]},
        {"label": "유형 2", "title_l1": "예측/최적화", "title_l2": "→ 운영비 절감",
         "cases": ["BP $1B 절감", "Tyson 적재율 46→87%", "조달 최적화 $287.9M"]},
        {"label": "유형 3", "title_l1": "오류/리스크 감소", "title_l2": "→ 손실 방지",
         "cases": ["배치 시간 83%↓", "재원 기간 1.5일 단축", "월 20병상 추가"]},
        {"label": "유형 4", "title_l1": "AIP Bootcamp", "title_l2": "5일 만에 성과",
         "cases": ["2일→$10M 절감", "5일→7개 ERP 통합", "첫 해 $100M 예상"]},
    ]

    for i, t in enumerate(types):
        x = card_x(n, i)
        add_rect(slide, x, card_y, card_w, card_h,
                 fill_color=C_BOX_BG, border_color=C_DARK_BG)

        add_textbox(slide, x, card_y + SP_XS, card_w, Inches(0.22),
                    text=t["label"], font_size=10, bold=True,
                    color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)

        inner_x = x + SP_MD
        inner_w = card_w - SP_MD * 2
        add_multiline_textbox(
            slide, inner_x, card_y + Inches(0.30), inner_w, Inches(0.60),
            [
                {"text": t["title_l1"], "font_size": 14, "bold": True, "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT},
                {"text": t["title_l2"], "font_size": 14, "bold": True, "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT},
            ]
        )

        add_rect(slide, inner_x, card_y + Inches(0.95), inner_w, Pt(0.5), fill_color=C_DARK_BG)

        case_lines = [{"text": c, "font_size": 11, "bold": False,
                       "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT} for c in t["cases"]]
        add_multiline_textbox(slide, inner_x, card_y + Inches(1.05),
                              inner_w, Inches(1.30), case_lines)

    # 하단 시사점
    bar = add_rect(slide, MARGIN_LEFT, Inches(4.35), CONTENT_W, Inches(0.40),
                   fill_color=C_BOX_BG, border_color=C_DARK_BG)
    txBox = slide.shapes.add_textbox(Inches(0.85), Inches(4.35), Inches(8.30), Inches(0.40))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "day1company 시사점: "
    r1.font.size = Pt(12)
    r1.font.bold = True
    r1.font.color.rgb = C_TEXT_BLACK
    r1.font.name = FONT_FAMILY
    r2 = p.add_run()
    r2.text = "유형 1(수작업 자동화)이 Layer 1 첫 프로젝트에 직접 대응"
    r2.font.size = Pt(12)
    r2.font.color.rgb = C_TEXT_BLACK
    r2.font.name = FONT_FAMILY


def slide_11_l2_cases(prs):
    """슬라이드 11: Layer 2 참고 사례 (화이트, 4열 카드 + 하단 매핑)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "Layer 2의 핵심은 기술이 아니라 제도 — 이 기업들은 제도를 먼저 설계", chapter="Part 2 — AI 전환 전략 프레임워크")

    n = 4
    card_w = PRESETS[n]["w"]
    card_h = Inches(1.60)
    card_y = ZONE_CONTENT_Y

    companies = [
        {"name": "Klarna", "policy": "채용 동결 +\n보상 재설계",
         "results": ["인당 매출 73%↑", "급여 60%↑"]},
        {"name": "Shopify", "policy": "평가 체계에\nAI 활용 내장",
         "results": ["인당 매출 127%↑", "매출 91%↑"]},
        {"name": "Block", "policy": "CEO 주도\n전사 구조 전환",
         "results": ["40% 감축", "$600M 절감"]},
        {"name": "Duolingo", "policy": "외주 전환 +\n내부 생산성 집중",
         "results": ["해고 0명", "생산량 4~5x"]},
    ]

    for i, co in enumerate(companies):
        x = card_x(n, i)
        add_rect(slide, x, card_y, card_w, card_h,
                 fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_textbox(slide, x, card_y + SP_XS, card_w, Inches(0.28),
                    text=co["name"], font_size=14, bold=True,
                    color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK,
                    align=PP_ALIGN.CENTER)
        add_multiline_textbox(
            slide, x + SP_SM, card_y + Inches(0.38), card_w - SP_SM * 2, Inches(0.45),
            [{"text": line, "font_size": 11, "bold": False, "color": C_TEXT_BLACK,
              "align": PP_ALIGN.CENTER} for line in co["policy"].split("\n")]
        )
        add_rect(slide, x + SP_LG, card_y + Inches(0.88), card_w - SP_LG * 2, Pt(0.5), fill_color=C_DARK_BG)
        result_lines = [{"text": r, "font_size": 12, "bold": True, "color": C_TEXT_BLACK,
                         "align": PP_ALIGN.CENTER} for r in co["results"]]
        add_multiline_textbox(slide, x, card_y + Inches(0.95), card_w, Inches(0.55), result_lines)

    # 하단 매핑 섹션
    add_textbox(slide, MARGIN_LEFT, Inches(3.55), CONTENT_W, Inches(0.30),
                text="글로벌 기업 제도 → day1company 대응 제도", font_size=15, bold=True,
                color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK)

    mappings = [
        ("채용 기준 변경 (Klarna, Shopify, Block)", "AI-First 채용 원칙"),
        ("평가 체계 연동 (Shopify)", "생산성 KPI 의무화 + 동일 직무 경쟁 평가제"),
        ("보상 재설계 (Klarna)", "AI 챔피언 제도"),
        ("CEO 직접 주도 (Block, Amazon)", "CAIO + C-Level 정기 리뷰"),
    ]
    map_y = Inches(3.90)
    for left_text, right_text in mappings:
        add_textbox(slide, MARGIN_LEFT, map_y, Inches(4.20), Inches(0.26),
                    text=left_text, font_size=11, color=C_TEXT_BLACK)
        add_textbox(slide, Inches(4.90), map_y, Inches(0.30), Inches(0.26),
                    text="→", font_size=12, bold=True, color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(5.30), map_y, Inches(4.00), Inches(0.26),
                    text=right_text, font_size=11, bold=True, color=C_TEXT_BLACK)
        map_y += Inches(0.30)


def slide_12_roadmap_overview(prs):
    """슬라이드 12: 실행 로드맵 개요 (화이트, 3단계 타임라인)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "Layer 1은 2개월 안에 완성하고, Layer 2는 즉시 시작한다", chapter="Part 2 — AI 전환 전략 프레임워크")

    n = 3
    col_w = Inches(2.60)
    col_gap = Inches(0.10)

    stages = [
        {"label": "1단계", "period": "Month 1~2", "keyword": "빠른 완성",
         "l1": "매출 정산 → 강사료 → 광고비\n3개 영역 완성",
         "l2": "4가지 제도 동시 도입\n(채용·KPI·경쟁평가·챔피언)"},
        {"label": "2단계", "period": "Month 3~6", "keyword": "고도화 + 확장",
         "l1": "피드백 반영 고도화\n추가 영역(변동비, 손익) 확장",
         "l2": "부서별 AI 활용 사례 중\n표준화 가능 건 → L1 편입"},
        {"label": "3단계", "period": "Month 7~", "keyword": "정착",
         "l1": "전사 데이터 기반 구축\nAI 의사결정 보조",
         "l2": "AI가 일상 업무 방식으로 정착\n자생적 운영"},
    ]

    # content_width 내 3등분 계산
    col_w = Inches(2.60)
    total_3col = col_w * 3 + col_gap * 2  # 8.00"
    col_start = MARGIN_LEFT + (CONTENT_W - total_3col) / 2  # 중앙 정렬

    header_y = ZONE_CONTENT_Y
    header_h = Inches(0.40)
    l1_y = Inches(2.18)
    l2_y = Inches(3.25)
    block_h = Inches(0.95)

    for i, stage in enumerate(stages):
        x = col_start + i * (col_w + col_gap)

        # 헤더 — 단일 텍스트로 통합
        hdr = add_rect(slide, x, header_y, col_w, header_h,
                       fill_color=C_DARK_BG, border_color=C_DARK_BG)
        add_text_to_shape(hdr, f"{stage['label']} | {stage['period']} — {stage['keyword']}",
                          font_size=10, bold=True, color=C_TEXT_WHITE)

        # L1 블록 — 라벨을 별도 다크 배지로, 콘텐츠와 분리
        l1_badge = add_rect(slide, x, l1_y, Inches(0.55), Inches(0.20),
                            fill_color=C_DARK_BG, border_color=C_DARK_BG)
        add_text_to_shape(l1_badge, "L1", font_size=8, bold=True, color=C_TEXT_WHITE)

        l1_box = add_rect(slide, x, l1_y + Inches(0.22), col_w, block_h - Inches(0.22),
                          fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_multiline_textbox(
            slide, x + SP_SM, l1_y + Inches(0.26), col_w - SP_SM * 2, Inches(0.60),
            [{"text": line, "font_size": 10, "bold": False, "color": C_TEXT_BLACK,
              "align": PP_ALIGN.LEFT} for line in stage["l1"].split("\n")]
        )

        # L2 블록
        l2_badge = add_rect(slide, x, l2_y, Inches(0.55), Inches(0.20),
                            fill_color=C_DARK_BG, border_color=C_DARK_BG)
        add_text_to_shape(l2_badge, "L2", font_size=8, bold=True, color=C_TEXT_WHITE)

        l2_box = add_rect(slide, x, l2_y + Inches(0.22), col_w, block_h - Inches(0.22),
                          fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_multiline_textbox(
            slide, x + SP_SM, l2_y + Inches(0.26), col_w - SP_SM * 2, Inches(0.60),
            [{"text": line, "font_size": 10, "bold": False, "color": C_TEXT_BLACK,
              "align": PP_ALIGN.LEFT} for line in stage["l2"].split("\n")]
        )

        # 화살표 (1→2, 2→3) — 간격 중앙
        if i < 2:
            arrow_x = x + col_w
            add_textbox(slide, arrow_x, Inches(2.50), col_gap, Inches(0.30),
                        text="→", font_size=14, bold=True, color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)

    # L2→L1 편입 라벨 — 2단계 열 중앙
    mid_x = col_start + (col_w + col_gap)
    add_textbox(slide, mid_x, l2_y - Inches(0.12), col_w, Inches(0.14),
                text="▲ L2 성과 → L1 편입", font_size=8, bold=True,
                color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)

    add_insight_line(slide, "승인 시 L1은 내일 착수, L2는 즉시 시행")


# ══════════════════════════════════════════════════
# Part 3: Layer 1 — 첫 프로젝트 선정과 실행 (슬라이드 13~18)
# ══════════════════════════════════════════════════

def slide_13_tech_approach(prs):
    """슬라이드 13: Layer 1의 기술 접근 (화이트, 상하 2단 구조도)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "검증된 아키텍처 위에 빠르게 구축, 결과를 바로 확인", chapter="Part 3 — Layer 1 첫 프로젝트와 실행")

    # 상단: Palantir 5계층
    add_textbox(slide, MARGIN_LEFT, ZONE_CONTENT_Y, CONTENT_W, Inches(0.25),
                text="Palantir Foundry 5계층 참조 모델", font_size=12, bold=True, color=C_TEXT_BLACK)

    layers_5 = ["1. 데이터\n수집", "2. 정제\n가공", "3. 비즈니스\n매핑", "4. 대시보드", "5. AI\n의사결정"]
    node_w = Inches(1.50)
    node_h = Inches(0.55)
    arrow_w = Inches(0.30)
    flow_y = Inches(1.95)
    x = MARGIN_LEFT
    for i, text in enumerate(layers_5):
        if i > 0:
            add_textbox(slide, x, flow_y, arrow_w, node_h,
                        text="→", font_size=14, bold=True, color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)
            x += arrow_w
        fill = C_DARK_BG if i in [1, 2] else C_BOX_BG
        tc = C_TEXT_WHITE if i in [1, 2] else C_TEXT_BLACK
        box = add_rect(slide, x, flow_y, node_w, node_h, fill_color=fill, border_color=C_DARK_BG)
        add_text_to_shape(box, text.replace("\n   ", "\n"), font_size=10, bold=(i in [1, 2]), color=tc)
        x += node_w

    add_textbox(slide, Inches(2.50), Inches(2.55), Inches(3.00), Inches(0.20),
                text="현재 day1company 위치: 계층 2~3", font_size=9, bold=True, color=C_ACCENT)

    # 하단: 데이터 3계층
    add_textbox(slide, MARGIN_LEFT, Inches(2.90), CONTENT_W, Inches(0.25),
                text="데이터 3계층 구조", font_size=12, bold=True, color=C_TEXT_BLACK)

    data_layers = ["A. Master Data", "B. Transaction", "C. Accounting"]
    block_w = Inches(2.60)
    block_h = Inches(0.50)
    block_y = Inches(3.20)
    x = MARGIN_LEFT
    for i, text in enumerate(data_layers):
        if i > 0:
            add_textbox(slide, x, block_y, arrow_w, block_h,
                        text="→", font_size=14, bold=True, color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)
            x += arrow_w
        box = add_rect(slide, x, block_y, block_w, block_h, fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_text_to_shape(box, text, font_size=12, bold=True, color=C_TEXT_BLACK)
        x += block_w

    add_textbox(slide, MARGIN_LEFT, Inches(3.75), Inches(5.00), Inches(0.20),
                text="Master Data를 첫 2주 안에 안정화 → 이후 프로젝트에서 재활용", font_size=10,
                color=C_SUB_WHITE)

    add_insight_line(slide, "검증된 아키텍처 위에 빠르게 구축")


def slide_14_candidate_map(prs):
    """슬라이드 14: 후보 선별 — 퍼널 시각화 (재디자인)
    개선: 3단 퍼널 (8개 전체 → 5개 탈락 → 3개 선정) + 선정 기준 4개"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "8개 후보 중 실질 후보 3개 선별", chapter="Part 3 — Layer 1 첫 프로젝트와 실행")

    # ── 퍼널 3단 ──
    funnel_x = MARGIN_LEFT
    funnel_stages = [
        {
            "label": "전체 후보 8개",
            "width": Inches(8.60),
            "fill": C_BOX_BG,
            "items": ["매출 정산", "비용 정산(강사료)", "광고비", "변동비", "손익 보고", "수강 현황", "B2B/B2G", "ERP 추출"],
        },
        {
            "label": "탈락 5개",
            "width": Inches(7.00),
            "fill": C_BOX_BG,
            "items": [
                ("변동비", "선행"),
                ("손익", "선행"),
                ("수강", "운영"),
                ("B2B", "매출"),
                ("ERP", "인프라"),
            ],
        },
        {
            "label": "실질 후보 3개",
            "width": Inches(5.00),
            "fill": C_DARK_BG,
            "items": ["A. 매출 정산", "C. 광고비 집계", "G. 인건비 정산"],
        },
    ]

    stage_h = Inches(0.55)
    stage_gap = Inches(0.10)
    sy = ZONE_CONTENT_Y

    for si, stage in enumerate(funnel_stages):
        # 중앙 정렬
        sx = MARGIN_LEFT + (CONTENT_W - stage["width"]) / 2

        box = add_rect(slide, sx, sy, stage["width"], stage_h,
                       fill_color=stage["fill"], border_color=C_DARK_BG)

        if si == 0:
            add_text_to_shape(box, " | ".join(stage["items"]),
                              font_size=9, color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)
        elif si == 1:
            items_text = "  |  ".join([f"{n} ({r})" for n, r in stage["items"]])
            add_text_to_shape(box, items_text,
                              font_size=8, color=C_SUB_WHITE, align=PP_ALIGN.CENTER)
        else:
            add_text_to_shape(box, "  |  ".join(stage["items"]),
                              font_size=13, bold=True, color=C_TEXT_WHITE, align=PP_ALIGN.CENTER)

        # 화살표 (단계 사이)
        if si < 2:
            add_textbox(slide, MARGIN_LEFT + CONTENT_W / 2 - Inches(0.15),
                        sy + stage_h, Inches(0.30), stage_gap,
                        text="▼", font_size=10, color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)

        sy += stage_h + stage_gap

    # ── 선정 기준 4개 (하단 가로 배치) — 퍼널 아래에 충분한 간격 ──
    crit_label_y = sy + SP_SM
    criteria = ["비용 절감 가시성", "독립 구축 가능성", "경영진 체감도", "후속 확장 기반"]
    crit_w = Inches(2.05)
    crit_h = Inches(0.38)
    crit_y = crit_label_y + Inches(0.28)

    add_textbox(slide, MARGIN_LEFT, crit_label_y, CONTENT_W, Inches(0.25),
                text="선정 기준", font_size=12, bold=True,
                color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK)

    for ci, crit in enumerate(criteria):
        cx = MARGIN_LEFT + ci * (crit_w + SP_SM)
        box = add_rect(slide, cx, crit_y, crit_w, crit_h,
                       fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_text_to_shape(box, crit, font_size=11, bold=True, color=C_TEXT_BLACK)

    add_insight_line(slide, "비용 절감 + 독립 구축 가능 3개 선별")


def slide_15_comparison(prs):
    """슬라이드 15: 후보 비교표 + CAIO 권고 (화이트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "가장 빠르고 확실한 성과가 기대되는 영역", chapter="Part 3 — Layer 1 첫 프로젝트와 실행")

    # 표 데이터
    headers = ["기준", "A. 매출 정산", "C. 광고비 집계", "G. 인건비 정산"]
    rows = [
        ["수작업 시간", "월 ~32시간 (4일)", "조사 필요", "조사 필요"],
        ["비용 절감 가시성", "★★★", "조사 후 확정", "조사 후 확정"],
        ["독립 구축 가능성", "★★★", "★★", "★★"],
        ["경영진 체감도", "★★★", "★★~★★★", "★★"],
        ["후속 확장 기반", "★★★", "★★", "★★"],
        ["Palantir 유형", "유형 1", "유형 1+2 혼합", "유형 1"],
    ]

    table_x = MARGIN_LEFT
    table_y = ZONE_CONTENT_Y
    col_w = [Inches(1.80), Inches(2.20), Inches(2.20), Inches(2.20)]
    row_h = Inches(0.35)

    # 헤더
    for j, header in enumerate(headers):
        hx = table_x + sum(col_w[:j])
        fill = C_DARK_BG if j == 0 or j == 1 else C_BOX_BG
        tc = C_TEXT_WHITE if j == 0 or j == 1 else C_TEXT_BLACK
        box = add_rect(slide, hx, table_y, col_w[j], row_h, fill_color=fill, border_color=C_DARK_BG)
        add_text_to_shape(box, header, font_size=11, bold=True, color=tc)

    # 데이터 행
    for i, row in enumerate(rows):
        ry = table_y + (i + 1) * row_h
        for j, cell in enumerate(row):
            cx = table_x + sum(col_w[:j])
            fill = C_BOX_BG if j == 0 else (RGBColor(0xE8, 0xE8, 0xE8) if j == 1 else C_WHITE_BG)
            box = add_rect(slide, cx, ry, col_w[j], row_h, fill_color=fill, border_color=C_BOX_BG)
            add_text_to_shape(box, cell, font_size=10, bold=(j == 0), color=C_TEXT_BLACK)

    # CAIO 권고 배너
    banner_y = table_y + 7 * row_h + SP_LG
    banner = add_rect(slide, MARGIN_LEFT, banner_y, CONTENT_W, Inches(0.40),
                      fill_color=C_DARK_BG, border_color=C_DARK_BG)
    add_text_to_shape(banner, "CAIO 권고: 매출 정산 자동화 — 현재 확보된 데이터 기준으로 가장 빠르고 확실한 선택",
                      font_size=12, bold=True, color=C_TEXT_WHITE)

    add_insight_line(slide, "최종 결정은 경영진의 판단 사항")


def slide_16_revenue_problem(prs):
    """슬라이드 16: 매출 정산 문제 (재디자인)
    개선: 좌측에 Big Number + 프로세스 흐름(가로), 우측에 6가지 상품유형 그리드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "매출 정산을 선택할 경우, 이것이 해결해야 할 구체적인 문제", chapter="Part 3 — Layer 1 첫 프로젝트와 실행")

    # ── 상단: Big Number (전폭) ──
    big_box = add_rect(slide, MARGIN_LEFT, ZONE_CONTENT_Y, CONTENT_W, Inches(0.70),
                       fill_color=C_DARK_BG, border_color=C_DARK_BG)
    add_textbox(slide, MARGIN_LEFT, ZONE_CONTENT_Y + Inches(0.05), Inches(4.00), Inches(0.55),
                text="월 32시간 (4일)", font_size=28, bold=True,
                color=C_TEXT_WHITE, font_name=FONT_FAMILY_BLACK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.00), ZONE_CONTENT_Y + Inches(0.15), Inches(4.30), Inches(0.40),
                text="매출 정산에 매월 소요되는 시간\n비용 정산 3일(24시간)은 별도",
                font_size=11, color=C_SUB_DARK, align=PP_ALIGN.LEFT)

    # ── 중단 좌측: 현재 프로세스 (가로 흐름) ──
    add_textbox(slide, MARGIN_LEFT, Inches(2.55), Inches(4.20), Inches(0.22),
                text="현재 프로세스", font_size=12, bold=True, color=C_TEXT_BLACK)

    steps = ["ERP\n다운로드", "수작업\n계산", "수작업\n검증", "보고"]
    node_w = Inches(0.90)
    node_h = Inches(0.50)
    arrow_w = Inches(0.25)
    fx = MARGIN_LEFT
    flow_y = Inches(2.82)

    for j, step in enumerate(steps):
        if j > 0:
            add_textbox(slide, fx, flow_y, arrow_w, node_h,
                        text="→", font_size=14, bold=True, color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)
            fx += arrow_w
        box = add_rect(slide, fx, flow_y, node_w, node_h,
                       fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_text_to_shape(box, step, font_size=10, bold=True, color=C_TEXT_BLACK)
        fx += node_w

    # ── 중단 우측: 6가지 상품유형 (2행×3열 그리드) ──
    add_textbox(slide, Inches(5.30), Inches(2.55), Inches(4.00), Inches(0.22),
                text="상품유형별 매출 인식 방식 — 6가지가 모두 다름", font_size=12, bold=True, color=C_TEXT_BLACK)

    types = ["단건 수강권", "구독형", "B2B 계약", "환불/부분취소", "번들 상품", "제휴/입점"]
    type_x = Inches(5.30)
    type_gap = Inches(0.08)
    type_w = (Inches(9.30) - type_x - 2 * type_gap) / 3  # 3열이 9.30" 내에 수용
    type_h = Inches(0.38)
    type_y = Inches(2.82)

    for ti, t in enumerate(types):
        row = ti // 3
        col = ti % 3
        tx = type_x + col * (type_w + type_gap)
        ty = type_y + row * (type_h + type_gap)
        box = add_rect(slide, tx, ty, type_w, type_h,
                       fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_text_to_shape(box, t, font_size=10, bold=True, color=C_TEXT_BLACK)

    # ── 하단: 핵심 문제 요약 ──
    summary = add_rect(slide, MARGIN_LEFT, Inches(4.00), CONTENT_W, Inches(0.50),
                       fill_color=C_BOX_BG, border_color=C_DARK_BG)
    add_multiline_textbox(
        slide, MARGIN_LEFT + SP_LG, Inches(4.02), CONTENT_W - SP_LG * 2, Inches(0.46),
        [
            {"text": "규칙이 명확하지만 수작업으로 반복 — 6가지 상품유형 × 7개 사업부문 = 42가지 조합을 수동 처리",
             "font_size": 11, "bold": False, "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT},
            {"text": "Palantir 유형 1(수작업 자동화)과 동일한 구조",
             "font_size": 11, "bold": True, "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT},
        ]
    )

    add_insight_line(slide, "수작업 반복 — Palantir 유형 1과 동일")


def slide_17_execution_plan(prs):
    """슬라이드 17: 2개월 실행 계획 (재수정 — 캔버스 내 맞춤)
    수정: 라벨+4열이 content_width(8.60") 내에 수용되도록 열 폭 조정"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "2개월 안에 3개 영역을 완성한다", chapter="Part 3 — Layer 1 첫 프로젝트와 실행")

    # 레이아웃 계산: 라벨(0.90") + 4열 = content_width(8.60") 내
    label_w = Inches(0.90)
    label_gap = SP_XS  # 0.08"
    grid_x = MARGIN_LEFT + label_w + label_gap
    available_w = Inches(8.60) - label_w - label_gap  # 7.62"
    col_gap = SP_XS  # 0.08"
    col_w = (available_w - 3 * col_gap) / 4  # ~1.845"

    periods = ["Week 1~2", "Week 3~4", "Week 5~6", "Week 7~8"]
    header_y = ZONE_CONTENT_Y
    header_h = Inches(0.28)

    for i, period in enumerate(periods):
        x = grid_x + i * (col_w + col_gap)
        box = add_rect(slide, x, header_y, col_w, header_h,
                       fill_color=C_DARK_BG, border_color=C_DARK_BG)
        add_text_to_shape(box, period, font_size=10, bold=True, color=C_TEXT_WHITE)

    rows_data = [
        {"label": "매출 정산", "cells": ["첫 작동 데모", "주력 사업부문\n완성", "전체 완성\n+ 병행 운영", "고도화"]},
        {"label": "강사료", "cells": [None, None, "설계 + 완성", "고도화"]},
        {"label": "광고비", "cells": [None, None, None, "착수 + 첫 작동"]},
    ]

    row_h = Inches(0.55)
    row_gap = SP_XS

    for ri, row in enumerate(rows_data):
        ry = header_y + header_h + SP_SM + ri * (row_h + row_gap)

        add_textbox(slide, MARGIN_LEFT, ry, label_w, row_h,
                    text=row["label"], font_size=10, bold=True, color=C_TEXT_BLACK, align=PP_ALIGN.LEFT)

        for ci, cell in enumerate(row["cells"]):
            if cell is None:
                continue
            cx = grid_x + ci * (col_w + col_gap)
            fill = C_DARK_BG if (ri == 0 and ci == 0) else C_BOX_BG
            tc = C_TEXT_WHITE if (ri == 0 and ci == 0) else C_TEXT_BLACK
            box = add_rect(slide, cx, ry, col_w, row_h,
                           fill_color=fill, border_color=C_DARK_BG)
            add_text_to_shape(box, cell, font_size=10, bold=False, color=tc)

    add_insight_line(slide, "먼저 완성하고 결과를 보여준 뒤, 쓰면서 고도화")


def slide_18_success_metrics(prs):
    """슬라이드 18: 성공 지표 + 실패 모드 방어선 (화이트, 2분할)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "가장 큰 위험은 기술 실패가 아니라 결과에 대한 불신", chapter="Part 3 — Layer 1 첫 프로젝트와 실행")

    n = 2
    card_w = PRESETS[n]["w"]
    card_y = ZONE_CONTENT_Y

    # 좌측: 성공 지표
    lx = card_x(n, 0)
    add_textbox(slide, lx, card_y, card_w, Inches(0.25),
                text="성공 지표", font_size=14, bold=True, color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK)

    metrics = [
        ("속도", "정산 시간 4일 → 1시간"),
        ("정확도", "수작업 대사 일치율 99%+"),
        ("범위", "매출 금액 80%+"),
    ]
    my = card_y + Inches(0.35)
    for label, value in metrics:
        box = add_rect(slide, lx, my, card_w, Inches(0.55),
                       fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_textbox(slide, lx + SP_SM, my + Inches(0.02), card_w - SP_SM * 2, Inches(0.22),
                    text=label, font_size=12, bold=True, color=C_TEXT_BLACK)
        add_textbox(slide, lx + SP_SM, my + Inches(0.25), card_w - SP_SM * 2, Inches(0.25),
                    text=value, font_size=14, bold=True, color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK)
        my += Inches(0.65)

    # 우측: 실패 모드 방어선
    rx = card_x(n, 1)
    add_textbox(slide, rx, card_y, card_w, Inches(0.25),
                text="실패 모드 방어선", font_size=14, bold=True, color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK)

    risks = [
        ("불신 ★", "대사 자동화 + 불일치 원인\n투명 공개 + 담당자 판정 권한"),
        ("미사용", "첫 주부터 담당자 참여 +\nLegacy 접근 권한 정리"),
        ("지연", "2주 내 첫 작동 공개로\n모멘텀 확보"),
    ]
    ry = card_y + Inches(0.35)
    for label, defense in risks:
        fill = C_DARK_BG if "★" in label else C_BOX_BG
        tc = C_TEXT_WHITE if "★" in label else C_TEXT_BLACK
        box = add_rect(slide, rx, ry, card_w, Inches(0.55),
                       fill_color=fill, border_color=C_DARK_BG)
        add_textbox(slide, rx + SP_SM, ry + Inches(0.02), Inches(1.00), Inches(0.22),
                    text=label, font_size=11, bold=True, color=tc)
        add_multiline_textbox(
            slide, rx + SP_SM, ry + Inches(0.22), card_w - SP_SM * 2, Inches(0.30),
            [{"text": line, "font_size": 9, "bold": False, "color": tc, "align": PP_ALIGN.LEFT}
             for line in defense.split("\n")]
        )
        ry += Inches(0.65)

    add_insight_line(slide, "방어선 설계 + 빠른 결과 공개가 최선")


# ══════════════════════════════════════════════════
# Part 4: Layer 2 — 전사 AI 활용 제도 설계 (슬라이드 19~23)
# ══════════════════════════════════════════════════

def slide_19_why_policy(prs):
    """슬라이드 19: 왜 제도가 필요한가 (화이트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "기술 도입이 아니라 제도 설계가 Layer 2의 본질", chapter="Part 4 — Layer 2 제도 설계 (4가지)")

    obstacles = [
        {"num": "①", "problem": "유인 없음", "desc": "AI 활용이 평가에 미반영 → 변화 거부",
         "solution": "제도 1, 2, 3"},
        {"num": "②", "problem": "신뢰 부족", "desc": "외부 대상 산출물의 AI 품질 우려",
         "solution": "Layer 1 + 품질 기준"},
        {"num": "③", "problem": "유지보수 공백", "desc": "담당자 부재 시 기존 방식으로 회귀",
         "solution": "제도 4 (AI 챔피언)"},
    ]

    y = ZONE_CONTENT_Y
    for obs in obstacles:
        # 좌측 장애물
        prob_box = add_rect(slide, MARGIN_LEFT, y, Inches(3.80), Inches(0.70),
                            fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_textbox(slide, MARGIN_LEFT + SP_SM, y + Inches(0.02), Inches(0.30), Inches(0.25),
                    text=obs["num"], font_size=14, bold=True, color=C_TEXT_BLACK)
        add_textbox(slide, Inches(1.10), y + Inches(0.02), Inches(2.50), Inches(0.25),
                    text=obs["problem"], font_size=14, bold=True, color=C_TEXT_BLACK)
        add_textbox(slide, Inches(1.10), y + Inches(0.30), Inches(3.00), Inches(0.35),
                    text=obs["desc"], font_size=10, color=C_SUB_WHITE)

        # 화살표
        add_textbox(slide, Inches(4.70), y + Inches(0.15), Inches(0.40), Inches(0.40),
                    text="→", font_size=18, bold=True, color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)

        # 우측 해결
        sol_box = add_rect(slide, Inches(5.20), y, Inches(4.10), Inches(0.70),
                           fill_color=C_DARK_BG, border_color=C_DARK_BG)
        add_text_to_shape(sol_box, obs["solution"], font_size=14, bold=True, color=C_TEXT_WHITE)

        y += Inches(0.85)

    add_insight_line(slide, "기술 도입이 아니라 제도 설계가 Layer 2의 본질")


def slide_20_policy1(prs):
    """슬라이드 20: 제도 1 — AI-First 채용 원칙 (화이트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "사람을 더 뽑는 것이 아니라, AI를 먼저 활용하는 것이 기본값이 된다", chapter="Part 4 [1/4] 진입 차단 — AI-First 채용 원칙")

    # 의사결정 플로우
    flow_items = [
        ("인력 충원\n요청 발생", C_BOX_BG, C_TEXT_BLACK),
        ("AI로\n대체 가능?", C_DARK_BG, C_TEXT_WHITE),
        ("AI 활용\n방안 수립", C_DARK_BG, C_TEXT_WHITE),
        ("AI-Augmented\nJD 작성", C_BOX_BG, C_TEXT_BLACK),
    ]

    node_w = Inches(1.60)
    node_h = Inches(0.60)
    x = MARGIN_LEFT
    flow_y = Inches(1.80)
    for i, (text, fill, tc) in enumerate(flow_items):
        if i > 0:
            add_textbox(slide, x, flow_y, Inches(0.35), node_h,
                        text="→", font_size=14, bold=True, color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)
            x += Inches(0.35)
        box = add_rect(slide, x, flow_y, node_w, node_h, fill_color=fill, border_color=C_DARK_BG)
        add_text_to_shape(box, text, font_size=11, bold=True, color=tc)
        x += node_w

    # 예시 박스
    ex_box = add_rect(slide, MARGIN_LEFT, Inches(2.60), CONTENT_W, Inches(0.50),
                      fill_color=C_BOX_BG, border_color=C_DARK_BG)
    add_text_to_shape(ex_box,
        "예시: 5명 팀 → 1명 추가 요청 → AI 활용으로 인당 1.4개 런칭 소명",
        font_size=11, color=C_TEXT_BLACK)

    # 구체 내용
    details = [
        "채용: 단순 결원 충원 동결. 채용 요청 시 AI 활용 가능성 소명",
        "외주: AI 대체 가능 영역 외주 예산 재검토. 절감액 일부 팀 환원",
        "참고: Shopify \"AI로 불가능한 이유를 먼저 증명\", Klarna 자연 퇴직 시 AI 흡수",
    ]
    dy = Inches(3.30)
    for d in details:
        add_textbox(slide, MARGIN_LEFT + SP_LG, dy, CONTENT_W - SP_LG, Inches(0.22),
                    text=d, font_size=10, color=C_TEXT_BLACK)
        dy += Inches(0.28)

    add_insight_line(slide, "AI 먼저 활용이 기본값")


def slide_21_policy2(prs):
    """슬라이드 21: 제도 2 — 생산성 KPI 의무화 (화이트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "모든 조직장에게 동일한 생산성 목표를 부여한다", chapter="Part 4 [2/4] 목표 부여 — 인당 생산성 KPI")

    # 피라미드 3단
    levels = [
        {"text": "조직장 KPI — 인당 생산성 향상 목표 의무, 미달 시 고과 상위 등급 제한",
         "w": Inches(5.50), "fill": C_DARK_BG, "tc": C_TEXT_WHITE},
        {"text": "조직 KPI — 미달 조직 → 조직 효율화 검토",
         "w": Inches(7.00), "fill": C_BOX_BG, "tc": C_TEXT_BLACK},
        {"text": "인센티브 — 최대 개선 부서 → 인센티브 지급",
         "w": Inches(8.40), "fill": C_BOX_BG, "tc": C_TEXT_BLACK},
    ]

    py = Inches(1.80)
    level_h = Inches(0.50)
    level_gap = SP_XS
    for lvl in levels:
        lx = MARGIN_LEFT + (CONTENT_W - lvl["w"]) / 2
        box = add_rect(slide, lx, py, lvl["w"], level_h,
                       fill_color=lvl["fill"], border_color=C_DARK_BG)
        add_text_to_shape(box, lvl["text"], font_size=11, bold=True, color=lvl["tc"])
        py += level_h + level_gap

    # 설계 원리 — 피라미드 아래에 충분한 간격
    principle_y = py + SP_LG
    principle_box = add_rect(slide, MARGIN_LEFT, principle_y, CONTENT_W, Inches(0.35),
                             fill_color=C_BOX_BG, border_color=C_DARK_BG)
    add_text_to_shape(principle_box,
        '설계 원리: "안 해도 되는 것" → "하지 않으면 불이익"',
        font_size=12, bold=True, color=C_TEXT_BLACK)

    add_insight_line(slide, "AI 활용이 선택이 아닌 필수")


def slide_22_policy3(prs):
    """슬라이드 22: 제도 3 — 동일 직무 경쟁 평가제 (화이트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "같은 일을 하는 팀 중 하나가 AI로 앞서면, 나머지는 따라와야 한다", chapter="Part 4 [3/4] 경쟁 유도 — 동일 직무 경쟁 평가제")

    # VS 구조
    left_box = add_rect(slide, Inches(0.90), Inches(1.80), Inches(3.50), Inches(1.20),
                        fill_color=C_DARK_BG, border_color=C_DARK_BG)
    add_multiline_textbox(
        slide, Inches(1.00), Inches(1.85), Inches(3.30), Inches(1.10),
        [
            {"text": "선도 부서", "font_size": 10, "bold": True, "color": C_SUB_DARK, "align": PP_ALIGN.LEFT},
            {"text": "A사업부 마케팅팀", "font_size": 14, "bold": True, "color": C_TEXT_WHITE, "align": PP_ALIGN.LEFT},
            {"text": "AI로 캠페인 10→12건 (20%↑)", "font_size": 11, "bold": False, "color": C_TEXT_WHITE, "align": PP_ALIGN.LEFT},
            {"text": "→ 고과 S 확정", "font_size": 12, "bold": True, "color": C_TEXT_WHITE, "align": PP_ALIGN.LEFT},
        ]
    )

    add_textbox(slide, Inches(4.50), Inches(2.10), Inches(1.00), Inches(0.50),
                text="VS", font_size=20, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    right_box = add_rect(slide, Inches(5.60), Inches(1.80), Inches(3.50), Inches(1.20),
                         fill_color=C_BOX_BG, border_color=C_DARK_BG)
    add_multiline_textbox(
        slide, Inches(5.70), Inches(1.85), Inches(3.30), Inches(1.10),
        [
            {"text": "대상 부서", "font_size": 10, "bold": True, "color": C_SUB_WHITE, "align": PP_ALIGN.LEFT},
            {"text": "B사업부 마케팅팀", "font_size": 14, "bold": True, "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT},
            {"text": "20% 미달성", "font_size": 11, "bold": False, "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT},
            {"text": "→ A 이상 불가", "font_size": 12, "bold": True, "color": C_ACCENT, "align": PP_ALIGN.LEFT},
        ]
    )

    # 측정 기준
    add_textbox(slide, MARGIN_LEFT, Inches(3.20), CONTENT_W, Inches(0.25),
                text="측정 기준", font_size=12, bold=True, color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK)
    criteria = ["월말 P&L 기준으로 확인", "3개월 연속 달성 시 인정", "부서별 사전 정의 지표 (인당 처리량, 인당 산출물 수 등)"]
    cy = Inches(3.50)
    for c in criteria:
        add_textbox(slide, MARGIN_LEFT + SP_LG, cy, CONTENT_W - SP_LG, Inches(0.22),
                    text=c, font_size=11, color=C_TEXT_BLACK)
        cy += Inches(0.26)

    add_insight_line(slide, "AI로 앞선 팀이 기준이 되는 구조")


def slide_23_policy4(prs):
    """슬라이드 23: 제도 4 — AI 챔피언 제도 (화이트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "챔피언은 부서의 변화를 이끄는 동시에, Layer 1과의 연결점이 된다", chapter="Part 4 [4/4] 지속 운영 — AI 챔피언 제도")

    # 중앙 허브
    hub = add_rect(slide, Inches(4.00), Inches(2.20), Inches(2.00), Inches(0.60),
                   fill_color=C_DARK_BG, border_color=C_DARK_BG)
    add_text_to_shape(hub, "AI 챔피언", font_size=16, bold=True, color=C_TEXT_WHITE)

    # 4방향 스포크
    spokes = [
        (Inches(1.00), Inches(1.80), "① 부서 내 선도"),
        (Inches(7.00), Inches(1.80), "② 동료 지원"),
        (Inches(1.00), Inches(3.10), "③ 사례 발굴"),
        (Inches(7.00), Inches(3.10), "④ Layer 1 연결"),
    ]
    for sx, sy, text in spokes:
        box = add_rect(slide, sx, sy, Inches(2.00), Inches(0.45),
                       fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_text_to_shape(box, text, font_size=11, bold=True, color=C_TEXT_BLACK)

    # 보상
    add_textbox(slide, Inches(3.50), Inches(1.70), Inches(3.00), Inches(0.25),
                text="보상: 성과 평가 가점 + 별도 인센티브", font_size=10, bold=True, color=C_TEXT_BLACK,
                align=PP_ALIGN.CENTER)

    # L2→L1 편입 경로
    add_rect(slide, MARGIN_LEFT, Inches(3.85), CONTENT_W, Inches(0.35),
             fill_color=C_BOX_BG, border_color=C_DARK_BG)
    add_textbox(slide, MARGIN_LEFT + SP_SM, Inches(3.85), CONTENT_W - SP_SM * 2, Inches(0.35),
                text="챔피언이 발굴 → CAIO에 전달 → L1 편입 검토", font_size=11, bold=True,
                color=C_TEXT_BLACK, align=PP_ALIGN.CENTER)

    add_insight_line(slide, "챔피언이 변화를 이끌고 L1과 연결")


# ══════════════════════════════════════════════════
# Part 5: 통합 로드맵과 기대 효과 (슬라이드 24~26)
# ══════════════════════════════════════════════════

def slide_24_integrated_timeline(prs):
    """슬라이드 24: 통합 타임라인 (슬라이드 12와 유사하지만 더 상세)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "Layer 1은 2개월 안에 완성하고, Layer 2는 즉시 시작한다", chapter="Part 5 — 통합 로드맵과 기대 효과")

    # 4구간 타임라인
    periods = [
        {"label": "즉시", "l1": "—", "l2": "AI-First 채용 원칙\n조직장 KPI 반영"},
        {"label": "Month 1~2", "l1": "매출 정산→강사료\n→광고비 완성", "l2": "AI 챔피언 1차 선정\n동일 직무 그룹 지정"},
        {"label": "Month 3~6", "l1": "고도화 + 추가\n영역 확장", "l2": "L2 성과 → L1 편입\n분기 리뷰"},
        {"label": "Month 7~", "l1": "전사 데이터 기반\nAI 의사결정 보조", "l2": "AI 일상화\n자생적 운영"},
    ]

    col_w = Inches(2.00)
    col_gap = SP_SM
    header_h = Inches(0.30)
    block_h = Inches(0.80)
    header_y = ZONE_CONTENT_Y
    l1_y = Inches(2.05)
    l2_y = Inches(3.00)

    for i, p in enumerate(periods):
        x = MARGIN_LEFT + i * (col_w + col_gap)

        hdr = add_rect(slide, x, header_y, col_w, header_h, fill_color=C_DARK_BG, border_color=C_DARK_BG)
        add_text_to_shape(hdr, p["label"], font_size=10, bold=True, color=C_TEXT_WHITE)

        l1_box = add_rect(slide, x, l1_y, col_w, block_h, fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_textbox(slide, x + SP_XS, l1_y + Inches(0.02), Inches(0.60), Inches(0.18),
                    text="L1", font_size=8, bold=True, color=C_TEXT_BLACK)
        add_multiline_textbox(
            slide, x + SP_XS, l1_y + Inches(0.18), col_w - SP_XS * 2, Inches(0.55),
            [{"text": line, "font_size": 10, "bold": False, "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT}
             for line in p["l1"].split("\n")]
        )

        l2_box = add_rect(slide, x, l2_y, col_w, block_h, fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_textbox(slide, x + SP_XS, l2_y + Inches(0.02), Inches(0.60), Inches(0.18),
                    text="L2", font_size=8, bold=True, color=C_TEXT_BLACK)
        add_multiline_textbox(
            slide, x + SP_XS, l2_y + Inches(0.18), col_w - SP_XS * 2, Inches(0.55),
            [{"text": line, "font_size": 10, "bold": False, "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT}
             for line in p["l2"].split("\n")]
        )

    add_insight_line(slide, "두 Layer가 합류하면서 상호 강화")


def slide_25_expected_effects(prs):
    """슬라이드 25: 기대 효과 (화이트, 2열 + 하단 배너)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "Layer 1의 시간 절약 + Layer 2의 채용 동결 = 실질적 인건비 절감", chapter="Part 5 — 통합 로드맵과 기대 효과")

    n = 2
    card_w = PRESETS[n]["w"]
    card_y = ZONE_CONTENT_Y
    card_h = Inches(1.80)

    # 좌측: L1 직접 효과
    lx = card_x(n, 0)
    add_textbox(slide, lx, card_y, card_w, Inches(0.25),
                text="Layer 1 직접 효과 (측정 가능)", font_size=13, bold=True, color=C_TEXT_BLACK)
    l1_items = [
        "매출 정산: 32h → 수시간",
        "강사료 정산: 조사 후 확정",
        "광고비 정산: 조사 후 확정",
    ]
    iy = card_y + Inches(0.35)
    for item in l1_items:
        box = add_rect(slide, lx, iy, card_w, Inches(0.35), fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_text_to_shape(box, item, font_size=12, color=C_TEXT_BLACK, align=PP_ALIGN.LEFT)
        iy += Inches(0.42)

    # 우측: L2 간접 효과
    rx = card_x(n, 1)
    add_textbox(slide, rx, card_y, card_w, Inches(0.25),
                text="Layer 2 간접 효과 (환산)", font_size=13, bold=True, color=C_TEXT_BLACK)
    l2_items = [
        "채용 동결 → 인건비 자연 감소",
        "외주비 절감 → 직접 비용 절감",
        "부서별 생산성↑ → 동일 인원 더 많은 성과",
    ]
    iy = card_y + Inches(0.35)
    for item in l2_items:
        box = add_rect(slide, rx, iy, card_w, Inches(0.35), fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_text_to_shape(box, item, font_size=12, color=C_TEXT_BLACK, align=PP_ALIGN.LEFT)
        iy += Inches(0.42)

    # 하단 배너
    banner = add_rect(slide, MARGIN_LEFT, Inches(3.80), CONTENT_W, Inches(0.70),
                      fill_color=C_DARK_BG, border_color=C_DARK_BG)
    add_textbox(slide, MARGIN_LEFT, Inches(3.82), CONTENT_W, Inches(0.40),
                text="인당 매출 2.34억 → 3억 (28%↑)", font_size=22, bold=True,
                color=C_TEXT_WHITE, font_name=FONT_FAMILY_BLACK, align=PP_ALIGN.CENTER)
    add_textbox(slide, MARGIN_LEFT, Inches(4.22), CONTENT_W, Inches(0.25),
                text="Klarna 73%↑ / Shopify 127%↑ 대비 보수적 목표", font_size=10,
                color=C_SUB_DARK, align=PP_ALIGN.CENTER)

    add_insight_line(slide, "인당 매출 3억 원은 달성 가능한 목표")


def slide_26_caio_role(prs):
    """슬라이드 26: CAIO 역할 체계 (화이트, 2분할)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "CAIO가 구조를 설계하고, 조직이 스스로 성장하는 구조", chapter="Part 5 — 통합 로드맵과 기대 효과")

    # 좌측: 역할
    add_textbox(slide, MARGIN_LEFT, ZONE_CONTENT_Y, Inches(4.00), Inches(0.25),
                text="CAIO 역할", font_size=13, bold=True, color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK)

    roles = [
        ("Layer 1", "직접 설계·구축\n(온톨로지, 자동화 파이프라인)"),
        ("Layer 2", "제도 설계 + 실행 지원\n(4가지 제도 운영, 챔피언 육성)"),
    ]
    ry = Inches(2.00)
    for label, desc in roles:
        box = add_rect(slide, MARGIN_LEFT, ry, Inches(4.00), Inches(0.80),
                       fill_color=C_BOX_BG, border_color=C_DARK_BG)
        add_textbox(slide, MARGIN_LEFT + SP_SM, ry + Inches(0.05), Inches(1.00), Inches(0.22),
                    text=label, font_size=12, bold=True, color=C_TEXT_BLACK)
        add_multiline_textbox(
            slide, MARGIN_LEFT + SP_SM, ry + Inches(0.28), Inches(3.70), Inches(0.45),
            [{"text": line, "font_size": 11, "bold": False, "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT}
             for line in desc.split("\n")]
        )
        ry += Inches(0.90)

    # 우측: 성장 경로
    add_textbox(slide, Inches(5.30), ZONE_CONTENT_Y, Inches(4.00), Inches(0.25),
                text="조직 성장 경로", font_size=13, bold=True, color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK)

    phases = [
        ("Phase 1", "CAIO 1인 — L1 구축 + L2 설계"),
        ("Phase 2", "챔피언 합류 + AI 챔피언 확산"),
        ("Phase 3", "내부 팀 + L2→L1 편입 정례화"),
        ("Phase 4", "AI 일상화, 인당 매출 3억 원"),
    ]
    py = Inches(2.00)
    for label, desc in phases:
        box = add_rect(slide, Inches(5.30), py, Inches(4.00), Inches(0.42),
                       fill_color=C_BOX_BG, border_color=C_DARK_BG)
        txBox = slide.shapes.add_textbox(Inches(5.40), py, Inches(3.80), Inches(0.42))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{label}  "
        r1.font.size = Pt(11)
        r1.font.bold = True
        r1.font.color.rgb = C_TEXT_BLACK
        r1.font.name = FONT_FAMILY
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = C_TEXT_BLACK
        r2.font.name = FONT_FAMILY
        py += Inches(0.50)

    add_insight_line(slide, "CAIO가 구조를 설계, 조직이 성장")


# ══════════════════════════════════════════════════
# Part 6: 예상 리스크와 대응 설계 (슬라이드 27~29)
# ══════════════════════════════════════════════════

def _risk_table_slide(prs, title, period_label, period_color, risks, insight_text, chapter="Part 6 — 예상 리스크와 대응"):
    """리스크 표 슬라이드 공통 함수"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, title, chapter=chapter)

    # 시간축 라벨
    add_textbox(slide, MARGIN_LEFT, ZONE_CONTENT_Y, Inches(1.50), Inches(0.25),
                text=period_label, font_size=12, bold=True, color=period_color)

    # 표 헤더
    headers = ["리스크", "대상", "설명", "대응 방안"]
    col_ws = [Inches(1.80), Inches(0.60), Inches(3.00), Inches(3.20)]
    header_y = ZONE_CONTENT_Y + Inches(0.30)
    row_h = Inches(0.30)

    hx = MARGIN_LEFT
    for j, (header, cw) in enumerate(zip(headers, col_ws)):
        box = add_rect(slide, hx, header_y, cw, row_h, fill_color=C_DARK_BG, border_color=C_DARK_BG)
        add_text_to_shape(box, header, font_size=10, bold=True, color=C_TEXT_WHITE)
        hx += cw

    # 데이터 행
    data_row_h = Inches(0.38)
    for i, risk in enumerate(risks):
        ry = header_y + row_h + i * data_row_h
        hx = MARGIN_LEFT
        for j, (cell, cw) in enumerate(zip(risk, col_ws)):
            fill = C_BOX_BG if i % 2 == 0 else C_WHITE_BG
            box = add_rect(slide, hx, ry, cw, data_row_h, fill_color=fill, border_color=C_BOX_BG)
            b = j == 0
            add_text_to_shape(box, cell, font_size=9, bold=b, color=C_TEXT_BLACK, align=PP_ALIGN.LEFT)
            hx += cw

    add_insight_line(slide, insight_text)


def slide_27_risk_early(prs):
    """슬라이드 27: 도입 초기 리스크"""
    risks = [
        ["L1 결과 불신", "L1", "자동화 결과 신뢰 부족", "담당자 첫 주 참여 + 대사 투명 공개"],
        ['"우리 부서는 다르다"', "L2", "현장 특수성으로 AI 거부", "CAIO 직접 참여, 범위 조정"],
        ["조직장 반발", "L2", "채용 동결·KPI 반영 이의", "사전 면담 + 글로벌 사례 공유"],
        ["L1 구축 지연", "L1", "데이터/ERP 연동 난이도", "2주 내 첫 작동 결과 공개"],
    ]
    _risk_table_slide(prs,
        "초기 리스크의 핵심은 '불신'과 '거부'이다",
        "Month 1~2", C_ACCENT, risks,
        "불신과 거부 — 빠른 결과 공개로 대응")


def slide_28_risk_mid(prs):
    """슬라이드 28: 운영 안정기 리스크"""
    risks = [
        ["제도의 형식화", "L2", "KPI만 맞추고 실질 AI 미활용", "P&L 기반 실질 산출물 측정"],
        ["조직장 방관", "L2", "제도 전달 안 하고 묵인", "타 부서 20%↑ 시 고과 상한"],
        ["보상 불만", "L2", "일은 늘고 보상은 동일", "외주비 절감 환원 + 챔피언 인센티브"],
        ["L2 성과 편차", "L2", "부서별 양극화", "경쟁 평가제 + CAIO 직접 지원"],
        ["L1 우선순위 충돌", "L1", "동시 자동화 요청 분산", "P&L 기준 비용 절감 기대값 순"],
    ]
    _risk_table_slide(prs,
        "형식화와 양극화 — 동일 직무 경쟁 평가제가 대응한다",
        "Month 3~6", C_TEXT_BLACK, risks,
        "형식화와 양극화 — 경쟁 평가제로 대응")


def slide_29_risk_long(prs):
    """슬라이드 29: 장기 리스크"""
    risks = [
        ["AI 챔피언 이탈", "L2", "퇴사 시 AI 활용 후퇴", "L1 편입 + 복수 챔피언 지정"],
        ["L2→L1 편입 지연", "통합", "개인 노하우로 남음", "분기별 편입 검토 정례화"],
        ['"충분히 했다" 증후군', "L2", "현 수준에 안주", "외부 사례 참고 + 기준선 상향"],
        ["CAIO 1인 의존", "L1", "부재 시 유지보수 불가", "Phase 2 챔피언 육성 + 문서화"],
    ]
    _risk_table_slide(prs,
        "사람 의존과 안주 — L2→L1 편입 구조가 방지한다",
        "Month 7~", RGBColor(0x33, 0x99, 0x33), risks,
        "사람 의존과 안주 — L2→L1 편입으로 방지")


# ══════════════════════════════════════════════════
# Part 7: 의사결정 요청 (슬라이드 30~32)
# ══════════════════════════════════════════════════

def slide_30_summary(prs):
    """슬라이드 30: 핵심 요약 (다크)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK_BG)

    add_textbox(slide, MARGIN_LEFT, Inches(0.40), CONTENT_W, Inches(0.60),
                text="핵심 요약", font_size=26, bold=True,
                color=C_TEXT_WHITE, font_name=FONT_FAMILY_BLACK)

    # L1 요약
    add_multiline_textbox(
        slide, MARGIN_LEFT, Inches(1.20), Inches(4.00), Inches(1.50),
        [
            {"text": "Layer 1 — 구조화 자동화 (Palantir Approach)", "font_size": 14, "bold": True, "color": C_TEXT_WHITE, "align": PP_ALIGN.LEFT},
            {"text": "CAIO가 매출 정산 → 강사료 → 광고비", "font_size": 12, "bold": False, "color": C_SUB_DARK, "align": PP_ALIGN.LEFT},
            {"text": "2개월 내 완성", "font_size": 12, "bold": False, "color": C_SUB_DARK, "align": PP_ALIGN.LEFT},
            {"text": "이후 Layer 2와 함께 고도화", "font_size": 12, "bold": False, "color": C_SUB_DARK, "align": PP_ALIGN.LEFT},
        ]
    )

    # L2 요약
    add_multiline_textbox(
        slide, Inches(5.30), Inches(1.20), Inches(4.00), Inches(1.50),
        [
            {"text": "Layer 2 — AI 도구 활용", "font_size": 14, "bold": True, "color": C_TEXT_WHITE, "align": PP_ALIGN.LEFT},
            {"text": "4가지 제도로 부서 자율 채택 유도", "font_size": 12, "bold": False, "color": C_SUB_DARK, "align": PP_ALIGN.LEFT},
            {"text": "AI-First 채용 / KPI / 경쟁 평가 / 챔피언", "font_size": 12, "bold": False, "color": C_SUB_DARK, "align": PP_ALIGN.LEFT},
            {"text": "즉시 시행", "font_size": 12, "bold": False, "color": C_SUB_DARK, "align": PP_ALIGN.LEFT},
        ]
    )

    # 상호보완
    add_textbox(slide, MARGIN_LEFT, Inches(2.90), CONTENT_W, Inches(0.30),
                text="Top-down(CAIO) + Bottom-up(각 부서) = 전사 양방향 변화 | L2 성과 → L1 편입 → 시스템 표준화",
                font_size=11, color=C_SUB_DARK, align=PP_ALIGN.CENTER)

    # 큰 수치
    add_textbox(slide, MARGIN_LEFT, Inches(3.50), CONTENT_W, Inches(0.60),
                text="인당 매출 2.34억 → 3억 원 (28%↑)", font_size=26, bold=True,
                color=C_TEXT_WHITE, font_name=FONT_FAMILY_BLACK, align=PP_ALIGN.CENTER)

    add_textbox(slide, MARGIN_LEFT, Inches(4.20), CONTENT_W, Inches(0.25),
                text="Klarna 73%↑ / Shopify 127%↑ 대비 보수적 목표", font_size=10,
                color=C_SUB_DARK, align=PP_ALIGN.CENTER)


def slide_31_decision(prs):
    """슬라이드 31: 경영진 의사결정 요청 (화이트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE_BG)
    add_white_title(slide, "이후 실행 플랜", chapter="Part 7 — 의사결정 요청")

    n = 2
    card_w = PRESETS[n]["w"]
    card_y = ZONE_CONTENT_Y
    card_h = Inches(2.20)

    # 좌측: L1
    lx = card_x(n, 0)
    l1_box = add_rect(slide, lx, card_y, card_w, card_h,
                      fill_color=C_BOX_BG, border_color=C_DARK_BG)
    add_textbox(slide, lx + SP_SM, card_y + SP_XS, card_w - SP_SM * 2, Inches(0.25),
                text="Layer 1", font_size=14, bold=True, color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK)

    l1_items = [
        "① 첫 프로젝트 선정 승인\n   (매출 정산 권고, 또는 다른 후보)",
        "② 2개월 실행 계획 승인\n   (매출 정산 → 강사료 → 광고비)",
        "③ 담당자 프로젝트 참여 시간 확보",
    ]
    iy = card_y + Inches(0.40)
    for item in l1_items:
        add_multiline_textbox(
            slide, lx + SP_LG, iy, card_w - SP_LG * 2, Inches(0.50),
            [{"text": line, "font_size": 11, "bold": False, "color": C_TEXT_BLACK, "align": PP_ALIGN.LEFT}
             for line in item.split("\n")]
        )
        iy += Inches(0.55)

    # 우측: L2
    rx = card_x(n, 1)
    l2_box = add_rect(slide, rx, card_y, card_w, card_h,
                      fill_color=C_BOX_BG, border_color=C_DARK_BG)
    add_textbox(slide, rx + SP_SM, card_y + SP_XS, card_w - SP_SM * 2, Inches(0.25),
                text="Layer 2", font_size=14, bold=True, color=C_TEXT_BLACK, font_name=FONT_FAMILY_BLACK)

    l2_items = [
        "④ AI-First 채용 원칙 도입 승인",
        "⑤ 조직장 KPI에 인당 생산성 반영 승인",
        "⑥ 동일 직무 경쟁 평가제 도입 승인",
    ]
    iy = card_y + Inches(0.40)
    for item in l2_items:
        add_textbox(slide, rx + SP_LG, iy, card_w - SP_LG * 2, Inches(0.35),
                    text=item, font_size=11, color=C_TEXT_BLACK)
        iy += Inches(0.55)

    # 하단 메시지
    banner = add_rect(slide, MARGIN_LEFT, Inches(4.20), CONTENT_W, Inches(0.40),
                      fill_color=C_DARK_BG, border_color=C_DARK_BG)
    add_text_to_shape(banner,
        "Layer 1은 2개월 안에 결과를 보여준다 — Layer 2는 제도로 조직이 스스로 움직이게 한다",
        font_size=12, bold=True, color=C_TEXT_WHITE)

    add_insight_line(slide, "오늘 승인 시, 내일부터 실행")


def slide_32_closing(prs):
    """슬라이드 32: 끝 (다크)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK_BG)

    add_textbox(slide, MARGIN_LEFT, Inches(1.80), CONTENT_W, Inches(1.00),
                text="감사합니다", font_size=34, bold=True,
                color=C_TEXT_WHITE, font_name=FONT_FAMILY_BLACK, align=PP_ALIGN.CENTER)

    add_textbox(slide, MARGIN_LEFT, Inches(2.80), CONTENT_W, Inches(0.60),
                text="Q&A", font_size=22, color=C_SUB_DARK, align=PP_ALIGN.CENTER)

    add_multiline_textbox(
        slide, MARGIN_LEFT, Inches(4.20), CONTENT_W, Inches(0.80),
        [
            {"text": "부록 A. 이해관계자별 실행 전략 | B. Lock-in 프로세스 | C. 저항 대응 매뉴얼",
             "font_size": 10, "bold": False, "color": C_SUB_DARK, "align": PP_ALIGN.CENTER},
            {"text": "부록 D. 후보 상세 평가표 | E. K-IFRS 15 준수",
             "font_size": 10, "bold": False, "color": C_SUB_DARK, "align": PP_ALIGN.CENTER},
        ]
    )


# ── 메인 실행 ──────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = CANVAS_W
    prs.slide_height = CANVAS_H

    # Part 1: AI 시대, 비용 구조가 바뀌고 있다
    slide_01_cover(prs)
    slide_02_global_cases(prs)
    slide_03_new_normal(prs)
    slide_04_agenda(prs)

    # ── 중간 제목: Part 1 후반 (위기감) ──
    add_section_title(prs,
        "AI 시대, 비용 구조가 바뀌고 있다",
        "Part 1 — 우리에게도 해당되는 변화")

    slide_05_disruption(prs)
    slide_06_current_state(prs)
    slide_07_before_after(prs)

    # ── 중간 제목: Part 2 (전략 프레임워크) ──
    add_section_title(prs,
        "AI 전환 전략 프레임워크",
        "Part 2 — 두 가지 Layer로 구성된 전략")

    slide_08_two_layers(prs)
    slide_09_synergy(prs)

    slide_10_palantir(prs)
    slide_11_l2_cases(prs)
    slide_12_roadmap_overview(prs)

    # ── 중간 제목: Part 3 (Layer 1 상세) ──
    add_section_title(prs,
        "Layer 1 — Palantir Approach, 첫 프로젝트와 실행",
        "Part 3 — 무엇을, 언제까지, 어떻게")

    slide_13_tech_approach(prs)
    slide_14_candidate_map(prs)
    slide_15_comparison(prs)
    slide_16_revenue_problem(prs)
    slide_17_execution_plan(prs)
    slide_18_success_metrics(prs)

    # ── 중간 제목: Part 4 (Layer 2 제도) ──
    add_section_title(prs,
        "Layer 2 — AI 도구 활용 제도 설계",
        "Part 4 — 진입 차단 → 목표 부여 → 경쟁 유도 → 지속 운영")

    slide_19_why_policy(prs)
    slide_20_policy1(prs)
    slide_21_policy2(prs)
    slide_22_policy3(prs)
    slide_23_policy4(prs)

    # ── 중간 제목: Part 5 (통합) ──
    add_section_title(prs,
        "통합 로드맵과 기대 효과",
        "Part 5 — 두 Layer를 합치면")

    slide_24_integrated_timeline(prs)
    slide_25_expected_effects(prs)
    slide_26_caio_role(prs)

    # ── 중간 제목: Part 6 (리스크) ──
    add_section_title(prs,
        "예상 리스크와 대응 설계",
        "Part 6 — 초기 → 안정기 → 장기")

    slide_27_risk_early(prs)
    slide_28_risk_mid(prs)
    slide_29_risk_long(prs)

    # Part 7: 의사결정 요청 + Q&A
    slide_30_summary(prs)
    slide_31_decision(prs)
    slide_32_closing(prs)

    total = len(prs.slides)
    output_path = "/Users/kangmin/cowork/slide_diagram_guide/output_full.pptx"
    prs.save(output_path)
    print(f"Full deck ({total} slides) saved to: {output_path}")


if __name__ == "__main__":
    main()
